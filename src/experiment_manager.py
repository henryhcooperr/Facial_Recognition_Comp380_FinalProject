#!/usr/bin/env python3

import os
import sys
import json
import uuid
import time
import logging
import datetime
import shutil
import yaml
import pandas as pd
import numpy as np
import torch
import matplotlib.pyplot as plt
from pathlib import Path
from typing import Dict, List, Tuple, Optional, Union, Any
from enum import Enum
import optuna
import torch.nn as nn
from sklearn.metrics import confusion_matrix, roc_curve, auc
from sklearn.manifold import TSNE
import torch.nn.functional as F

from .base_config import PROJECT_ROOT, CHECKPOINTS_DIR, OUT_DIR, logger
from .data_prep import PreprocessingConfig, process_raw_data
from .face_models import get_model, get_criterion
from .testing import evaluate_model, plot_confusion_matrix, plot_roc_curves, plot_gradcam_visualization
from .visualize import plot_tsne_embeddings, plot_attention_maps, plot_embedding_similarity, plot_learning_curves
from .training_utils import save_checkpoint, prune_checkpoints, EarlyStopping, get_scheduler, apply_gradient_clipping, plot_lr_schedule
from .advanced_metrics import create_enhanced_confusion_matrix, plot_advanced_confusion_matrix, calculate_per_class_metrics, plot_class_difficulty_analysis, expected_calibration_error, plot_reliability_diagram, plot_confidence_histogram, plot_resource_usage


class ExperimentConfig:
    """Configuration for face recognition experiments."""
    
    class Dataset(Enum):
        DATASET1 = "dataset1"  # High diversity of subjects (36 people, 49 images each)
        DATASET2 = "dataset2"  # High quantity per person (18 people, 100 images each)
        BOTH = "both"          # Both datasets
    
    class ModelArchitecture(Enum):
        BASELINE = "baseline"
        CNN = "cnn"  # ResNet transfer learning
        SIAMESE = "siamese"
        ATTENTION = "attention"
        ARCFACE = "arcface"
        HYBRID = "hybrid"
    
    class LRSchedulerType(Enum):
        """Types of learning rate schedulers."""
        NONE = "none"
        STEP = "step"
        EXPONENTIAL = "exponential"
        COSINE = "cosine"
        REDUCE_ON_PLATEAU = "reduce_on_plateau"
        ONE_CYCLE = "one_cycle"
    
    class EvaluationMode(Enum):
        """Evaluation modes for model evaluation."""
        STANDARD = "standard"  # Standard accuracy, precision, recall, F1
        ENHANCED = "enhanced"  # Standard + per-class metrics, calibration, resource usage
        
    def __init__(
        self,
        experiment_name: str = "Unnamed Experiment",
        dataset: Union[Dataset, str] = Dataset.BOTH,
        model_architecture: Union[ModelArchitecture, str, List[str]] = ModelArchitecture.CNN,
        preprocessing_config: Optional[PreprocessingConfig] = None,
        epochs: int = 30,
        batch_size: int = 32,
        learning_rate: float = 0.001,
        cross_dataset_testing: bool = True,
        results_dir: Optional[str] = None,
        experiment_id: Optional[str] = None,
        config_version: str = "1.0.0",
        config_history: Optional[List[Dict[str, Any]]] = None,
        
        # Training enhancements
        random_seed: int = 42,
        use_early_stopping: bool = False,
        early_stopping_patience: int = 10,
        early_stopping_min_delta: float = 0.0,
        early_stopping_metric: str = "loss",
        early_stopping_mode: str = "min",
        use_gradient_clipping: bool = False,
        gradient_clipping_max_norm: float = 1.0,
        gradient_clipping_adaptive: bool = False,
        lr_scheduler_type: Union[LRSchedulerType, str] = LRSchedulerType.NONE,
        lr_scheduler_params: Optional[Dict[str, Any]] = None,
        
        # Enhanced model checkpointing parameters (pre-existing)
        save_best_checkpoint: bool = True,
        checkpoint_frequency: int = 1,
        keep_last_n_checkpoints: int = 5,
        keep_best_n_checkpoints: int = 3,
        save_checkpoint_metadata: bool = True,
        resumable_training: bool = True,
        
        # Enhanced evaluation metrics (new)
        evaluation_mode: Union[EvaluationMode, str] = EvaluationMode.STANDARD,
        per_class_analysis: bool = True,
        calibration_analysis: bool = True,
        resource_monitoring: bool = True,
        calibration_n_bins: int = 10,
        use_temperature_scaling: bool = False,
        confidence_threshold: float = 0.5,
        save_raw_predictions: bool = True,
        max_misclassified_examples: int = 10,
        model_complexity_analysis: bool = True,
        class_difficulty_metric: str = "f1"
    ):
        """Initialize experiment configuration.
        
        Args:
            experiment_name: Human-readable name for the experiment
            dataset: Which dataset to use for training/testing
            model_architecture: Model architecture to use 
            preprocessing_config: Config for preprocessing steps
            epochs: Number of training epochs
            batch_size: Training batch size
            learning_rate: Initial learning rate
            cross_dataset_testing: Whether to test on other datasets
            results_dir: Directory to save results (auto-generated if None)
            experiment_id: Unique ID for experiment (auto-generated if None)
            config_version: Version of this configuration
            config_history: List of previous configurations
            
            # Training enhancements
            random_seed: Random seed for reproducibility
            use_early_stopping: Whether to use early stopping
            early_stopping_patience: Patience for early stopping
            early_stopping_min_delta: Minimum change to count as improvement
            early_stopping_mode: 'min' for loss, 'max' for accuracy
            early_stopping_metric: 'loss' or 'accuracy'
            use_gradient_clipping: Whether to use gradient clipping
            gradient_clipping_max_norm: Maximum gradient norm
            gradient_clipping_adaptive: Whether to use adaptive clipping
            lr_scheduler_type: Type of learning rate scheduler to use
            lr_scheduler_params: Parameters for the scheduler
            
            # Checkpointing
            save_best_checkpoint: Whether to save the best model checkpoint
            checkpoint_frequency: How often to save checkpoints (in epochs)
            keep_last_n_checkpoints: Number of recent checkpoints to keep
            keep_best_n_checkpoints: Number of best checkpoints to keep
            save_checkpoint_metadata: Whether to save metadata with checkpoints
            resumable_training: Whether to allow resuming from checkpoints
            
            # Enhanced evaluation metrics
            evaluation_mode: Mode for model evaluation
            per_class_analysis: Whether to perform per-class analysis
            calibration_analysis: Whether to analyze prediction calibration
            resource_monitoring: Whether to monitor resource usage
            calibration_n_bins: Number of bins for calibration metrics
            use_temperature_scaling: Whether to use temperature scaling
            confidence_threshold: Threshold for binary predictions
            save_raw_predictions: Whether to save raw predictions
            max_misclassified_examples: Max number of examples to store
            model_complexity_analysis: Whether to analyze model complexity
            class_difficulty_metric: Metric to use for class difficulty ranking
        """
        # Set basic configuration parameters
        self.experiment_name = experiment_name
        self.experiment_id = experiment_id or str(uuid.uuid4())[:8]
        
        # Creation timestamp
        self.created_at = datetime.datetime.now().isoformat()
        
        # Handle dataset parameter
        if isinstance(dataset, str):
            try:
                self.dataset = ExperimentConfig.Dataset(dataset)
            except ValueError:
                raise ValueError(f"Invalid dataset: {dataset}")
        else:
            self.dataset = dataset
        
        # Handle model_architecture parameter (could be enum, string, or list)
        if isinstance(model_architecture, list):
            # For comparison experiments with multiple architectures
            self.model_architecture = model_architecture
        elif isinstance(model_architecture, str):
            try:
                self.model_architecture = ExperimentConfig.ModelArchitecture(model_architecture)
            except ValueError:
                raise ValueError(f"Invalid model architecture: {model_architecture}")
        else:
            self.model_architecture = model_architecture
        
        # Set training parameters
        self.preprocessing_config = preprocessing_config
        self.epochs = epochs
        self.batch_size = batch_size
        self.learning_rate = learning_rate
        self.cross_dataset_testing = cross_dataset_testing
        
        # Set up results directory
        if results_dir:
            self.results_dir = Path(results_dir)
        else:
            # Auto-generate results directory name
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            self.results_dir = OUT_DIR / f"experiment_{self.experiment_id}_{timestamp}"
        
        # Configuration versioning
        self.config_version = config_version
        self.config_history = config_history or []
        
        # Training enhancements
        self.random_seed = random_seed
        self.use_early_stopping = use_early_stopping
        self.early_stopping_patience = early_stopping_patience
        self.early_stopping_min_delta = early_stopping_min_delta
        self.early_stopping_metric = early_stopping_metric
        self.early_stopping_mode = early_stopping_mode
        self.use_gradient_clipping = use_gradient_clipping
        self.gradient_clipping_max_norm = gradient_clipping_max_norm
        self.gradient_clipping_adaptive = gradient_clipping_adaptive
        
        # Handle LR scheduler type
        if isinstance(lr_scheduler_type, str):
            try:
                self.lr_scheduler_type = ExperimentConfig.LRSchedulerType(lr_scheduler_type)
            except ValueError:
                logger.warning(f"Invalid scheduler type: {lr_scheduler_type}. Using NONE.")
                self.lr_scheduler_type = ExperimentConfig.LRSchedulerType.NONE
        else:
            self.lr_scheduler_type = lr_scheduler_type
        
        self.lr_scheduler_params = lr_scheduler_params or {}
        
        # Enhanced model checkpointing parameters
        self.save_best_checkpoint = save_best_checkpoint
        self.checkpoint_frequency = checkpoint_frequency
        self.keep_last_n_checkpoints = keep_last_n_checkpoints
        self.keep_best_n_checkpoints = keep_best_n_checkpoints
        self.save_checkpoint_metadata = save_checkpoint_metadata
        self.resumable_training = resumable_training
        
        # Enhanced evaluation metrics
        if isinstance(evaluation_mode, str):
            try:
                self.evaluation_mode = ExperimentConfig.EvaluationMode(evaluation_mode)
            except ValueError:
                raise ValueError(f"Invalid evaluation mode: {evaluation_mode}")
        else:
            self.evaluation_mode = evaluation_mode
            
        self.per_class_analysis = per_class_analysis
        self.calibration_analysis = calibration_analysis
        self.resource_monitoring = resource_monitoring
        self.calibration_n_bins = calibration_n_bins
        self.use_temperature_scaling = use_temperature_scaling
        self.confidence_threshold = confidence_threshold
        self.save_raw_predictions = save_raw_predictions
        self.max_misclassified_examples = max_misclassified_examples
        self.model_complexity_analysis = model_complexity_analysis
        self.class_difficulty_metric = class_difficulty_metric
    
    def increment_version(self, level='patch'):
        """Increment the configuration version.
        
        Args:
            level: Which part of the version to increment ('major', 'minor', or 'patch')
        """
        # Add current version to history before incrementing
        self.add_to_history()
        
        # Parse current version
        major, minor, patch = map(int, self.config_version.split('.'))
        
        # Increment appropriate part
        if level == 'major':
            major += 1
            minor = 0
            patch = 0
        elif level == 'minor':
            minor += 1
            patch = 0
        else:  # patch
            patch += 1
            
        # Update version
        self.config_version = f"{major}.{minor}.{patch}"
        return self.config_version
    
    def add_to_history(self):
        """Add current configuration state to history."""
        # Create a copy of current config state
        current_state = self.to_dict()
        current_state.pop('config_history', None)  # Remove history to avoid nesting
        
        # Add to history with timestamp
        history_entry = {
            'timestamp': datetime.datetime.now().isoformat(),
            'config': current_state
        }
        self.config_history.append(history_entry)
    
    def compare_with(self, other: 'ExperimentConfig') -> Dict[str, Dict[str, Any]]:
        """Compare this config with another and return differences."""
        this_dict = self.to_dict()
        other_dict = other.to_dict()
        
        # Remove history and timestamps to focus on functional differences
        for d in [this_dict, other_dict]:
            d.pop('config_history', None)
            d.pop('created_at', None)
        
        # Find differences
        differences = {}
        all_keys = set(this_dict.keys()) | set(other_dict.keys())
        
        for key in all_keys:
            if key not in this_dict:
                differences[key] = {'status': 'added', 'value': other_dict[key]}
            elif key not in other_dict:
                differences[key] = {'status': 'removed', 'value': this_dict[key]}
            elif this_dict[key] != other_dict[key]:
                differences[key] = {
                    'status': 'modified',
                    'old_value': this_dict[key],
                    'new_value': other_dict[key]
                }
        
        return differences

    def to_dict(self) -> Dict:
        """Convert config to dictionary for saving."""
        config_dict = {
            "experiment_id": self.experiment_id,
            "experiment_name": self.experiment_name,
            "dataset": self.dataset.value,
            "model_architecture": self.model_architecture.value if not isinstance(self.model_architecture, list) else self.model_architecture,
            "epochs": self.epochs,
            "batch_size": self.batch_size,
            "learning_rate": self.learning_rate,
            "cross_dataset_testing": self.cross_dataset_testing,
            "results_dir": str(self.results_dir),
            "created_at": self.created_at,
            "random_seed": self.random_seed,
            "config_version": self.config_version,
            "config_history": self.config_history,
            "use_early_stopping": self.use_early_stopping,
            "early_stopping_patience": self.early_stopping_patience,
            "early_stopping_min_delta": self.early_stopping_min_delta,
            "early_stopping_metric": self.early_stopping_metric,
            "early_stopping_mode": self.early_stopping_mode,
            "use_gradient_clipping": self.use_gradient_clipping,
            "gradient_clipping_max_norm": self.gradient_clipping_max_norm,
            "gradient_clipping_adaptive": self.gradient_clipping_adaptive,
            "lr_scheduler_type": self.lr_scheduler_type.value if isinstance(self.lr_scheduler_type, Enum) else self.lr_scheduler_type,
            "lr_scheduler_params": self.lr_scheduler_params,
            "save_best_checkpoint": self.save_best_checkpoint,
            "checkpoint_frequency": self.checkpoint_frequency,
            "keep_last_n_checkpoints": self.keep_last_n_checkpoints,
            "keep_best_n_checkpoints": self.keep_best_n_checkpoints,
            "save_checkpoint_metadata": self.save_checkpoint_metadata,
            "resumable_training": self.resumable_training
        }
        
        # Add preprocessing config if available
        if self.preprocessing_config:
            config_dict["preprocessing_config"] = self.preprocessing_config.to_dict()
            
        return config_dict
    
    @classmethod
    def from_dict(cls, config_dict: Dict) -> 'ExperimentConfig':
        """Create config from dictionary."""
        # Extract preprocessing config if available
        preprocessing_config = None
        if "preprocessing_config" in config_dict:
            preprocessing_config = PreprocessingConfig.from_dict(config_dict["preprocessing_config"])
            
        # Create the basic config
        config = cls(
            experiment_id=config_dict.get("experiment_id"),
            experiment_name=config_dict.get("experiment_name", "Unnamed Experiment"),
            dataset=config_dict.get("dataset", "both"),
            model_architecture=config_dict.get("model_architecture", "cnn"),
            preprocessing_config=preprocessing_config,
            epochs=config_dict.get("epochs", 30),
            batch_size=config_dict.get("batch_size", 32),
            learning_rate=config_dict.get("learning_rate", 0.001),
            cross_dataset_testing=config_dict.get("cross_dataset_testing", True),
            results_dir=config_dict.get("results_dir"),
            random_seed=config_dict.get("random_seed", 42),
            config_version=config_dict.get("config_version", "1.0.0"),
            use_early_stopping=config_dict.get("use_early_stopping", False),
            early_stopping_patience=config_dict.get("early_stopping_patience", 10),
            early_stopping_min_delta=config_dict.get("early_stopping_min_delta", 0.0),
            early_stopping_metric=config_dict.get("early_stopping_metric", "loss"),
            early_stopping_mode=config_dict.get("early_stopping_mode", "min"),
            use_gradient_clipping=config_dict.get("use_gradient_clipping", False),
            gradient_clipping_max_norm=config_dict.get("gradient_clipping_max_norm", 1.0),
            gradient_clipping_adaptive=config_dict.get("gradient_clipping_adaptive", False),
            lr_scheduler_type=config_dict.get("lr_scheduler_type", "none"),
            lr_scheduler_params=config_dict.get("lr_scheduler_params", {}),
            save_best_checkpoint=config_dict.get("save_best_checkpoint", True),
            checkpoint_frequency=config_dict.get("checkpoint_frequency", 1),
            keep_last_n_checkpoints=config_dict.get("keep_last_n_checkpoints", 5),
            keep_best_n_checkpoints=config_dict.get("keep_best_n_checkpoints", 3),
            save_checkpoint_metadata=config_dict.get("save_checkpoint_metadata", True),
            resumable_training=config_dict.get("resumable_training", True)
        )
        
        # Set creation timestamp if available
        if "created_at" in config_dict:
            config.created_at = config_dict["created_at"]
            
        # Set config history if available
        if "config_history" in config_dict:
            config.config_history = config_dict["config_history"]
            
        return config
    
    def to_yaml(self) -> str:
        """Convert config to YAML string."""
        return yaml.dump(self.to_dict(), default_flow_style=False, sort_keys=False)
    
    @classmethod
    def from_yaml(cls, yaml_str: str) -> 'ExperimentConfig':
        """Create config from YAML string."""
        config_dict = yaml.safe_load(yaml_str)
        return cls.from_dict(config_dict)
    
    def save(self, filepath: Optional[Path] = None) -> Path:
        """Save config to JSON file."""
        if filepath is None:
            filepath = self.results_dir / "experiment_config.json"
            
        # Create directory if it doesn't exist
        filepath.parent.mkdir(parents=True, exist_ok=True)
        
        with open(filepath, 'w') as f:
            json.dump(self.to_dict(), f, indent=2)
            
        return filepath
    
    def save_yaml(self, filepath: Optional[Path] = None) -> Path:
        """Save config to YAML file."""
        if filepath is None:
            filepath = self.results_dir / "experiment_config.yaml"
            
        # Create directory if it doesn't exist
        filepath.parent.mkdir(parents=True, exist_ok=True)
        
        with open(filepath, 'w') as f:
            f.write(self.to_yaml())
            
        return filepath
    
    @classmethod
    def load(cls, filepath: Path) -> 'ExperimentConfig':
        """Load config from JSON file."""
        with open(filepath, 'r') as f:
            config_dict = json.load(f)
            
        return cls.from_dict(config_dict)
    
    @classmethod
    def load_yaml(cls, filepath: Path) -> 'ExperimentConfig':
        """Load config from YAML file."""
        with open(filepath, 'r') as f:
            yaml_str = f.read()
            
        return cls.from_yaml(yaml_str)


class ModelRegistry:
    """Registry to manage face recognition models."""
    
    def __init__(self, registry_path: Optional[Path] = None):
        """Initialize model registry."""
        self.registry_path = registry_path or CHECKPOINTS_DIR / "model_registry.json"
        self.models = {}
        
        # Load existing registry if it exists
        if self.registry_path.exists():
            self.load()
    
    def register_model(self, 
                      model_name: str, 
                      architecture: str,
                      dataset_name: str,
                      experiment_id: str,
                      parameters: Dict[str, Any],
                      metrics: Optional[Dict[str, float]] = None,
                      checkpoint_path: Optional[Path] = None) -> str:
        """Register a new model with metadata."""
        # Generate a timestamp
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        
        # Check if model name exists, add versioning if needed
        base_name = model_name
        version = 1
        
        while model_name in self.models:
            model_name = f"{base_name}_v{version}"
            version += 1
        
        # Create model metadata
        model_metadata = {
            "model_name": model_name,
            "architecture": architecture,
            "dataset": dataset_name,
            "experiment_id": experiment_id,
            "parameters": parameters,
            "created_at": timestamp,
            "metrics": metrics or {},
            "checkpoint_path": str(checkpoint_path) if checkpoint_path else None
        }
        
        # Add to registry
        self.models[model_name] = model_metadata
        
        # Save registry
        self.save()
        
        return model_name
    
    def get_model(self, model_name: str) -> Dict[str, Any]:
        """Get model metadata by name."""
        if model_name not in self.models:
            raise ValueError(f"Model not found: {model_name}")
        
        return self.models[model_name]
    
    def list_models(self, 
                   architecture: Optional[str] = None, 
                   dataset: Optional[str] = None) -> List[Dict[str, Any]]:
        """List models with optional filtering."""
        filtered_models = self.models.values()
        
        if architecture:
            filtered_models = [m for m in filtered_models if m["architecture"] == architecture]
            
        if dataset:
            filtered_models = [m for m in filtered_models if m["dataset"] == dataset]
            
        return list(filtered_models)
    
    def update_metrics(self, model_name: str, metrics: Dict[str, float]) -> None:
        """Update model metrics."""
        if model_name not in self.models:
            raise ValueError(f"Model not found: {model_name}")
        
        self.models[model_name]["metrics"].update(metrics)
        self.models[model_name]["updated_at"] = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        
        # Save registry
        self.save()
    
    def save(self) -> None:
        """Save registry to disk."""
        self.registry_path.parent.mkdir(parents=True, exist_ok=True)
        
        with open(self.registry_path, 'w') as f:
            json.dump(self.models, f, indent=2)
    
    def load(self) -> None:
        """Load registry from disk."""
        if not self.registry_path.exists():
            self.models = {}
            return
        
        with open(self.registry_path, 'r') as f:
            self.models = json.load(f)


class ResultsManager:
    """Manager for experiment results and metrics."""
    
    def __init__(self, experiment_config: ExperimentConfig):
        """Initialize results manager with experiment configuration."""
        self.config = experiment_config
        self.results_dir = experiment_config.results_dir
        
        # Create directory structure
        self._setup_directories()
        
        # Initialize metrics storage
        self.train_metrics = []
        self.eval_metrics = []
        self.test_metrics = []
        self.confusion_matrices = {}
        self.per_class_metrics = {}
        self.calibration_metrics = {}
        self.resource_metrics = {}
        self.experiment_log = []
        
        # Handle model_architecture correctly when it's a list
        model_arch = self.config.model_architecture
        model_arch_value = model_arch if isinstance(model_arch, list) else model_arch.value
        
        # Add experiment start entry
        self.log_event("experiment_started", {
            "experiment_id": self.config.experiment_id,
            "experiment_name": self.config.experiment_name,
            "dataset": self.config.dataset.value,
            "model_architecture": model_arch_value,
            "timestamp": datetime.datetime.now().isoformat()
        })
    
    def _setup_directories(self):
        """Create directory structure for results."""
        # Main directories
        self.results_dir.mkdir(parents=True, exist_ok=True)
        
        # Create subdirectories
        self.metrics_dir = self.results_dir / "metrics"
        self.metrics_dir.mkdir(exist_ok=True)
        
        self.plots_dir = self.results_dir / "plots"
        self.plots_dir.mkdir(exist_ok=True)
        
        # Create subdirectories for each type of visualization
        self.confusion_dir = self.plots_dir / "confusion_matrices"
        self.confusion_dir.mkdir(exist_ok=True)
        
        self.calibration_dir = self.plots_dir / "calibration"
        self.calibration_dir.mkdir(exist_ok=True)
        
        self.class_analysis_dir = self.plots_dir / "class_analysis"
        self.class_analysis_dir.mkdir(exist_ok=True)
        
        self.resource_dir = self.plots_dir / "resources"
        self.resource_dir.mkdir(exist_ok=True)
        
        self.checkpoints_dir = self.results_dir / "checkpoints"
        self.checkpoints_dir.mkdir(exist_ok=True)
        
        self.logs_dir = self.results_dir / "logs"
        self.logs_dir.mkdir(exist_ok=True)
        
        # New directory for raw predictions (for post-hoc analysis)
        self.predictions_dir = self.results_dir / "predictions"
        self.predictions_dir.mkdir(exist_ok=True)
    
    def record_training_metrics(self, epoch: int, metrics: Dict[str, float]):
        """Record per-epoch training metrics."""
        metrics = {**metrics, "epoch": epoch}
        self.train_metrics.append(metrics)
        
        # Log event
        self.log_event("training_epoch_completed", {
            "epoch": epoch,
            "metrics": metrics
        })
        
        # Save metrics to CSV
        self._save_metrics_to_csv(self.train_metrics, "training_metrics.csv")
    
    def record_evaluation_metrics(self, epoch: int, metrics: Dict[str, float], dataset: str = "validation"):
        """Record evaluation metrics."""
        metrics = {**metrics, "epoch": epoch, "dataset": dataset}
        self.eval_metrics.append(metrics)
        
        # Log event
        self.log_event("evaluation_completed", {
            "epoch": epoch,
            "dataset": dataset,
            "metrics": metrics
        })
        
        # Save metrics to CSV
        self._save_metrics_to_csv(self.eval_metrics, "evaluation_metrics.csv")
    
    def record_test_metrics(self, metrics: Dict[str, float], dataset: str = "test"):
        """Record test metrics after training is complete."""
        metrics = {**metrics, "dataset": dataset}
        self.test_metrics.append(metrics)
        
        # Log event
        self.log_event("testing_completed", {
            "dataset": dataset,
            "metrics": metrics
        })
        
        # Save metrics to CSV
        self._save_metrics_to_csv(self.test_metrics, "test_metrics.csv")
    
    def record_confusion_matrix(self, y_true: List[int], y_pred: List[int], 
                              class_names: List[str], dataset: str = "test"):
        """Record basic confusion matrix (backward compatibility)."""
        from .advanced_metrics import create_enhanced_confusion_matrix, plot_advanced_confusion_matrix
        
        # Use enhanced confusion matrix functionality
        enhanced_cm = create_enhanced_confusion_matrix(y_true, y_pred, class_names)
        
        # Store confusion matrix
        self.confusion_matrices[dataset] = enhanced_cm
        
        # Save to JSON
        cm_path = self.metrics_dir / f"confusion_matrix_{dataset}.json"
        with open(cm_path, 'w') as f:
            json.dump(enhanced_cm, f, indent=2)
        
        # Create enhanced visualization
        plot_advanced_confusion_matrix(
            enhanced_cm, 
            output_path=self.confusion_dir / f"confusion_matrix_{dataset}.png"
        )
    
    def record_per_class_metrics(self, y_true: List[int], y_pred: List[int], 
                                y_score: np.ndarray, class_names: List[str], 
                                dataset: str = "test"):
        """
        Record detailed per-class performance metrics.
        
        Args:
            y_true: Ground truth labels
            y_pred: Predicted labels
            y_score: Prediction confidence scores
            class_names: Names of the classes
            dataset: Dataset name (e.g., "test", "validation")
        """
        from .advanced_metrics import (
            calculate_per_class_metrics, 
            plot_class_difficulty_analysis
        )
        
        # Calculate per-class metrics
        per_class_metrics = calculate_per_class_metrics(
            np.array(y_true), np.array(y_pred), y_score, class_names
        )
        
        # Store per-class metrics
        self.per_class_metrics[dataset] = per_class_metrics
        
        # Save to JSON
        metrics_path = self.metrics_dir / f"per_class_metrics_{dataset}.json"
        with open(metrics_path, 'w') as f:
            json.dump(per_class_metrics, f, indent=2)
        
        # Create class difficulty visualization
        plot_class_difficulty_analysis(
            per_class_metrics, 
            output_path=self.class_analysis_dir / f"class_difficulty_{dataset}.png",
            metric="f1"  # Use F1 score for difficulty ranking
        )
        
        # Also create visualizations by precision and recall
        plot_class_difficulty_analysis(
            per_class_metrics, 
            output_path=self.class_analysis_dir / f"class_difficulty_precision_{dataset}.png",
            metric="precision"
        )
        
        plot_class_difficulty_analysis(
            per_class_metrics, 
            output_path=self.class_analysis_dir / f"class_difficulty_recall_{dataset}.png",
            metric="recall"
        )
        
        # Log event
        self.log_event("per_class_metrics_recorded", {
            "dataset": dataset,
            "num_classes": len(class_names)
        })
    
    def record_calibration_metrics(self, y_true: List[int], y_pred: List[int], 
                                  y_score: np.ndarray, dataset: str = "test"):
        """
        Record confidence calibration metrics.
        
        Args:
            y_true: Ground truth labels
            y_pred: Predicted labels
            y_score: Prediction confidence scores
            dataset: Dataset name (e.g., "test", "validation")
        """
        from .advanced_metrics import (
            expected_calibration_error, 
            plot_reliability_diagram,
            plot_confidence_histogram
        )
        
        # Calculate calibration metrics
        calibration_data = expected_calibration_error(
            np.array(y_true), np.array(y_pred), y_score
        )
        
        # Store calibration metrics
        self.calibration_metrics[dataset] = calibration_data
        
        # Save to JSON
        metrics_path = self.metrics_dir / f"calibration_metrics_{dataset}.json"
        with open(metrics_path, 'w') as f:
            json.dump(calibration_data, f, indent=2)
        
        # Create reliability diagram
        plot_reliability_diagram(
            calibration_data,
            output_path=self.calibration_dir / f"reliability_diagram_{dataset}.png"
        )
        
        # Create confidence histogram
        plot_confidence_histogram(
            np.array(y_true), np.array(y_pred), y_score,
            output_path=self.calibration_dir / f"confidence_histogram_{dataset}.png"
        )
        
        # Log event
        self.log_event("calibration_metrics_recorded", {
            "dataset": dataset,
            "ece": calibration_data["expected_calibration_error"]
        })
    
    def record_resource_metrics(self, resource_data: Dict[str, Any], phase: str):
        """
        Record resource utilization metrics.
        
        Args:
            resource_data: Dictionary with resource usage data
            phase: Phase of the experiment (e.g., "training", "inference")
        """
        from .advanced_metrics import plot_resource_usage
        
        # Store resource metrics
        self.resource_metrics[phase] = resource_data
        
        # Save to JSON
        metrics_path = self.metrics_dir / f"resource_metrics_{phase}.json"
        with open(metrics_path, 'w') as f:
            json.dump(resource_data, f, indent=2)
        
        # Create resource usage plot
        plot_resource_usage(
            resource_data,
            output_path=self.resource_dir / f"resource_usage_{phase}.png"
        )
        
        # Log event
        self.log_event("resource_metrics_recorded", {
            "phase": phase,
            "duration": resource_data["duration"],
            "avg_cpu": resource_data["cpu_percent"]["mean"],
            "avg_memory_mb": resource_data["memory_mb"]["mean"]
        })
    
    def record_model_complexity(self, model: torch.nn.Module, input_size: Tuple[int, ...]):
        """
        Record model complexity metrics (parameters, FLOPs).
        
        Args:
            model: PyTorch model
            input_size: Input size for the model (batch_size, channels, height, width)
        """
        from .advanced_metrics import count_model_parameters, estimate_model_flops
        
        # Count parameters
        param_counts = count_model_parameters(model)
        
        # Estimate FLOPs
        try:
            flop_estimates = estimate_model_flops(model, input_size)
        except Exception as e:
            logger.warning(f"Could not estimate FLOPs: {str(e)}")
            flop_estimates = {"error": str(e)}
        
        # Combine metrics
        complexity_metrics = {
            "parameters": param_counts,
            "flops_estimate": flop_estimates,
            "input_size": input_size
        }
        
        # Store complexity metrics
        self.resource_metrics["model_complexity"] = complexity_metrics
        
        # Save to JSON
        metrics_path = self.metrics_dir / "model_complexity.json"
        with open(metrics_path, 'w') as f:
            json.dump(complexity_metrics, f, indent=2)
        
        # Log event
        self.log_event("model_complexity_recorded", {
            "total_parameters": param_counts["total_parameters"],
            "trainable_parameters": param_counts["trainable_parameters"]
        })
    
    def save_raw_predictions(self, y_true: List[int], y_pred: List[int], 
                           y_score: np.ndarray, class_names: List[str], 
                           filenames: Optional[List[str]] = None, 
                           dataset: str = "test"):
        """
        Save raw predictions for post-hoc analysis.
        
        Args:
            y_true: Ground truth labels
            y_pred: Predicted labels
            y_score: Prediction confidence scores
            class_names: Names of the classes
            filenames: Optional list of filenames corresponding to the samples
            dataset: Dataset name (e.g., "test", "validation")
        """
        # Convert predictions to a DataFrame
        import pandas as pd
        
        # Create basic dataframe
        data = {
            "true_label_idx": y_true,
            "true_label": [class_names[i] for i in y_true],
            "pred_label_idx": y_pred,
            "pred_label": [class_names[i] for i in y_pred],
            "correct": np.array(y_true) == np.array(y_pred)
        }
        
        # Add confidence scores
        if y_score.ndim > 1:
            # For multi-class, get confidence for predicted class and all classes
            data["confidence"] = np.array([y_score[i, pred] for i, pred in enumerate(y_pred)])
            
            # Add confidence for each class
            for i, class_name in enumerate(class_names):
                if i < y_score.shape[1]:
                    data[f"conf_{class_name}"] = y_score[:, i]
        else:
            # For binary classification
            data["confidence"] = y_score
        
        # Add filenames if provided
        if filenames:
            data["filename"] = filenames
        
        # Create DataFrame
        df = pd.DataFrame(data)
        
        # Save to CSV
        predictions_path = self.predictions_dir / f"predictions_{dataset}.csv"
        df.to_csv(predictions_path, index=False)
        
        # Log event
        self.log_event("raw_predictions_saved", {
            "dataset": dataset,
            "num_samples": len(y_true),
            "accuracy": float(np.mean(np.array(y_true) == np.array(y_pred)))
        })
    
    def record_learning_curves(self, train_losses: List[float], val_losses: List[float], 
                             accuracies: List[float]):
        """Record and visualize learning curves."""
        # Save data
        curves_data = {
            "train_losses": train_losses,
            "val_losses": val_losses,
            "accuracies": accuracies
        }
        
        curves_path = self.metrics_dir / "learning_curves.json"
        with open(curves_path, 'w') as f:
            json.dump(curves_data, f, indent=2)
        
        # Create visualization
        plot_learning_curves(train_losses, val_losses, accuracies, 
                          str(self.plots_dir), self.config.experiment_id)
    
    def log_event(self, event_type: str, data: Dict[str, Any]):
        """Log an experiment event with timestamp."""
        log_entry = {
            "timestamp": datetime.datetime.now().isoformat(),
            "event_type": event_type,
            "data": data
        }
        
        self.experiment_log.append(log_entry)
        
        # Save log to file
        self._save_experiment_log()
    
    def _save_experiment_log(self):
        """Save experiment log to JSON file."""
        log_path = self.logs_dir / "experiment_log.json"
        with open(log_path, 'w') as f:
            json.dump(self.experiment_log, f, indent=2)
    
    def _save_metrics_to_csv(self, metrics_list: List[Dict[str, Any]], filename: str):
        """Save metrics to CSV file."""
        if not metrics_list:
            return
        
        df = pd.DataFrame(metrics_list)
        csv_path = self.metrics_dir / filename
        df.to_csv(csv_path, index=False)
    
    def save_model_checkpoint(self, model: nn.Module, optimizer: torch.optim.Optimizer, 
                            epoch: int, is_best: bool = False, 
                            scheduler: Optional[torch.optim.lr_scheduler._LRScheduler] = None,
                            metrics: Optional[Dict[str, float]] = None) -> Path:
        """Save model checkpoint with enhanced functionality.
        
        Args:
            model: PyTorch model
            optimizer: PyTorch optimizer
            epoch: Current epoch number
            is_best: Whether this is the best model so far
            scheduler: Learning rate scheduler (optional)
            metrics: Dictionary of validation metrics
            
        Returns:
            Path to saved checkpoint
        """
        # Create checkpoint filename
        checkpoint_filename = f"checkpoint_epoch_{epoch}.pth"
        best_model_path = self.checkpoints_dir / "best_model.pth"
        
        # Use metrics if provided or create minimal dict
        validation_metrics = metrics or {"is_best": is_best}
        
        # Create metadata
        metadata = {
            "experiment_id": self.config.experiment_id,
            "experiment_name": self.config.experiment_name,
            "model_architecture": self.config.model_architecture.value if hasattr(self.config.model_architecture, 'value') else self.config.model_architecture,
            "dataset": self.config.dataset.value,
            "preprocessing": self.config.preprocessing_config.name if self.config.preprocessing_config else "None",
            "config_version": self.config.config_version,
            "timestamp": datetime.datetime.now().isoformat()
        }
        
        # Save checkpoint
        checkpoint_path = save_checkpoint(
            model=model,
            optimizer=optimizer,
            scheduler=scheduler,
            epoch=epoch,
            validation_metrics=validation_metrics,
            checkpoint_dir=self.checkpoints_dir,
            filename=checkpoint_filename,
            metadata=metadata if self.config.save_checkpoint_metadata else None,
            keep_best_only=False  # We'll handle pruning separately
        )
        
        # Save best model if needed
        if is_best:
            # Save model state dict
            torch.save(model.state_dict(), best_model_path)
            
            # If full checkpoint for best model is requested
            if self.config.save_best_checkpoint:
                best_checkpoint_path = self.checkpoints_dir / "best_checkpoint.pth"
                # Save full checkpoint including optimizer, scheduler etc.
                save_checkpoint(
                    model=model,
                    optimizer=optimizer,
                    scheduler=scheduler,
                    epoch=epoch,
                    validation_metrics=validation_metrics,
                    checkpoint_dir=self.checkpoints_dir,
                    filename="best_checkpoint.pth",
                    metadata=metadata if self.config.save_checkpoint_metadata else None,
                    keep_best_only=False
                )
            
            # Log event
            self.log_event("best_model_saved", {"epoch": epoch})
        
        # Prune old checkpoints based on config
        if epoch % self.config.checkpoint_frequency == 0:
            # Keep the last N checkpoints
            if self.config.keep_last_n_checkpoints > 0:
                prune_checkpoints(
                    self.checkpoints_dir, 
                    keep=self.config.keep_last_n_checkpoints,
                    pattern="checkpoint_epoch_*.pth"
                )
        
        # Log event
        self.log_event("checkpoint_saved", {
            "epoch": epoch,
            "path": str(checkpoint_path),
            "is_best": is_best
        })
        
        return checkpoint_path if not is_best else best_model_path
    
    def generate_experiment_summary(self) -> Dict[str, Any]:
        """Generate a comprehensive experiment summary."""
        # Gather key metrics
        best_val_metrics = {}
        if self.eval_metrics:
            # Find epoch with best validation accuracy
            best_epoch_idx = max(range(len(self.eval_metrics)), 
                               key=lambda i: self.eval_metrics[i].get('accuracy', 0))
            best_val_metrics = self.eval_metrics[best_epoch_idx]
        
        # Compute average test metrics across datasets
        avg_test_metrics = {}
        if self.test_metrics:
            metrics_keys = set()
            for metrics in self.test_metrics:
                metrics_keys.update(metrics.keys())
            
            metrics_keys.discard('dataset')  # Remove non-numeric field
            
            # Calculate averages
            for key in metrics_keys:
                values = [m.get(key, 0) for m in self.test_metrics if key in m]
                if values:
                    avg_test_metrics[f"avg_{key}"] = sum(values) / len(values)
        
        # Handle model_architecture correctly when it's a list
        model_arch = self.config.model_architecture
        model_arch_value = model_arch if isinstance(model_arch, list) else model_arch.value
        
        # Extract training enhancements information
        training_enhancements = {
            "early_stopping": {
                "enabled": self.config.use_early_stopping,
                "patience": self.config.early_stopping_patience,
                "min_delta": self.config.early_stopping_min_delta,
                "metric": self.config.early_stopping_metric,
                "mode": self.config.early_stopping_mode
            },
            "gradient_clipping": {
                "enabled": self.config.use_gradient_clipping,
                "max_norm": self.config.gradient_clipping_max_norm,
                "adaptive": self.config.gradient_clipping_adaptive
            },
            "lr_scheduler": {
                "type": self.config.lr_scheduler_type.value if hasattr(self.config.lr_scheduler_type, 'value') else self.config.lr_scheduler_type,
                "params": self.config.lr_scheduler_params
            },
            "checkpointing": {
                "save_best": self.config.save_best_checkpoint,
                "frequency": self.config.checkpoint_frequency,
                "keep_last_n": self.config.keep_last_n_checkpoints,
                "keep_best_n": self.config.keep_best_n_checkpoints,
                "resumable": self.config.resumable_training
            }
        }
        
        # Create summary
        summary = {
            "experiment_id": self.config.experiment_id,
            "experiment_name": self.config.experiment_name,
            "dataset": self.config.dataset.value,
            "model_architecture": model_arch_value,
            "preprocessing": self.config.preprocessing_config.name if self.config.preprocessing_config else "None",
            "config_version": self.config.config_version,
            "training_params": {
                "epochs": self.config.epochs,
                "batch_size": self.config.batch_size,
                "learning_rate": self.config.learning_rate,
                "random_seed": self.config.random_seed
            },
            "training_enhancements": training_enhancements,
            "best_validation_metrics": best_val_metrics,
            "test_metrics": self.test_metrics,
            "average_test_metrics": avg_test_metrics,
            "completed_at": datetime.datetime.now().isoformat()
        }
        
        # Check if early stopping was triggered
        early_stopping_trace_path = self.logs_dir / "early_stopping_trace.json"
        if early_stopping_trace_path.exists():
            try:
                with open(early_stopping_trace_path, 'r') as f:
                    early_stopping_info = json.load(f)
                summary["early_stopping_info"] = early_stopping_info
            except Exception as e:
                logger.error(f"Error loading early stopping trace: {str(e)}")
        
        # Save summary to file
        summary_path = self.results_dir / "experiment_summary.json"
        with open(summary_path, 'w') as f:
            json.dump(summary, f, indent=2)
        
        # Log event
        self.log_event("experiment_completed", {
            "summary": summary
        })
        
        return summary
    
    def export_to_csv(self, filepath: Optional[Path] = None) -> Path:
        """Export all metrics to a single CSV file."""
        if filepath is None:
            filepath = self.results_dir / "all_metrics.csv"
        
        # Combine all metrics
        all_metrics = []
        
        # Add training metrics
        for metric in self.train_metrics:
            metric_copy = metric.copy()
            metric_copy['phase'] = 'training'
            all_metrics.append(metric_copy)
        
        # Add evaluation metrics
        for metric in self.eval_metrics:
            metric_copy = metric.copy()
            metric_copy['phase'] = 'evaluation'
            all_metrics.append(metric_copy)
        
        # Add test metrics
        for metric in self.test_metrics:
            metric_copy = metric.copy()
            metric_copy['phase'] = 'test'
            all_metrics.append(metric_copy)
        
        # Create DataFrame and save
        if all_metrics:
            df = pd.DataFrame(all_metrics)
            df.to_csv(filepath, index=False)
        
        return filepath 


class ExperimentManager:
    """Manager for running face recognition experiments."""
    
    def __init__(self, model_registry: Optional[ModelRegistry] = None):
        """Initialize experiment manager."""
        self.model_registry = model_registry or ModelRegistry()
    
    def load_experiment_config(self, config_path: Union[str, Path]) -> ExperimentConfig:
        """Load experiment configuration from file.
        
        Automatically detects file format based on extension (.json or .yaml/.yml).
        """
        config_path = Path(config_path)
        if not config_path.exists():
            raise FileNotFoundError(f"Configuration file not found: {config_path}")
            
        # Detect file format based on extension
        if config_path.suffix.lower() in ['.yaml', '.yml']:
            return ExperimentConfig.load_yaml(config_path)
        elif config_path.suffix.lower() == '.json':
            return ExperimentConfig.load(config_path)
        else:
            raise ValueError(f"Unsupported configuration file format: {config_path.suffix}")
        
    def run_experiment(self, config: Union[ExperimentConfig, str, Path]) -> Dict[str, Any]:
        """Run a face recognition experiment according to configuration.
        
        Args:
            config: Either an ExperimentConfig object or a path to a config file (.json/.yaml/.yml)
        """
        # If config is a string or Path, load it as a configuration file
        if isinstance(config, (str, Path)):
            config = self.load_experiment_config(config)
            
        # Set up results manager
        results_manager = ResultsManager(config)
        
        # Log experiment start
        logger.info(f"Starting experiment: {config.experiment_name} (ID: {config.experiment_id})")
        logger.info(f"Configuration: {config.to_dict()}")
        
        # Set random seeds for reproducibility
        from .base_config import set_random_seeds
        set_random_seeds(seed=config.random_seed)
        logger.info(f"Random seed set to: {config.random_seed}")
        
        # Save experiment configuration in both JSON and YAML formats for reproducibility
        config.save()
        config.save_yaml()
        
        # Determine experiment type based on configuration
        if config.dataset == ExperimentConfig.Dataset.BOTH and config.cross_dataset_testing:
            return self._run_cross_dataset_experiment(config, results_manager)
        elif isinstance(config.model_architecture, list):
            # If model_architecture is a list, run architecture comparison
            return self._run_architecture_comparison_experiment(config, results_manager)
        elif len(config.model_architecture.value) == 1:
            return self._run_single_model_experiment(config, results_manager)
        else:
            return self._run_architecture_comparison_experiment(config, results_manager)
    
    def _run_single_model_experiment(self, config: ExperimentConfig, 
                                   results_manager: ResultsManager) -> Dict[str, Any]:
        """Run an experiment with a single model architecture on one dataset."""
        try:
            # Set up device
            device = torch.device('cuda' if torch.cuda.is_available() else 'cpu')
            
            # Handle model_architecture correctly when it's a list
            model_arch = config.model_architecture
            model_arch_value = model_arch[0] if isinstance(model_arch, list) else model_arch.value
            
            # Import training utilities
            from .training_utils import (
                EarlyStopping, 
                get_scheduler, 
                apply_gradient_clipping,
                plot_lr_schedule
            )
            
            # Import resource monitoring utilities
            from .advanced_metrics import TimerContext, ResourceMonitor
            
            # Initialize model
            model = get_model(model_arch_value)
            model = model.to(device)
            
            # Record model complexity metrics
            input_shape = (config.batch_size, 3, 224, 224)  # Typical input shape for face recognition
            results_manager.record_model_complexity(model, input_shape)
            
            # Set up optimizer
            optimizer = torch.optim.Adam(model.parameters(), lr=config.learning_rate)
            
            # Get criterion based on model type
            criterion = get_criterion(model_arch_value)
            
            # Setup data loaders with seed for reproducibility
            train_loader, val_loader, test_loader, class_names = self._setup_data_loaders(
                config.dataset.value, 
                config.preprocessing_config, 
                batch_size=config.batch_size,
                seed=config.random_seed
            )
            
            # Set up learning rate scheduler if requested
            scheduler = None
            if config.lr_scheduler_type != ExperimentConfig.LRSchedulerType.NONE:
                scheduler_type = config.lr_scheduler_type.value
                # Calculate steps_per_epoch if needed for schedulers like OneCycleLR
                if scheduler_type == "one_cycle":
                    config.lr_scheduler_params['steps_per_epoch'] = len(train_loader)
                    config.lr_scheduler_params['epochs'] = config.epochs
                
                scheduler = get_scheduler(
                    scheduler_type=scheduler_type,
                    optimizer=optimizer,
                    **config.lr_scheduler_params
                )
                
                # Plot and save LR schedule if a scheduler is used
                if scheduler:
                    plot_lr_schedule(
                        scheduler=scheduler,
                        optimizer=optimizer,
                        num_epochs=config.epochs,
                        save_path=results_manager.plots_dir / "lr_schedule.png"
                    )
            
            # Set up early stopping if requested
            early_stopping = None
            if config.use_early_stopping:
                early_stopping = EarlyStopping(
                    patience=config.early_stopping_patience,
                    min_delta=config.early_stopping_min_delta,
                    mode=config.early_stopping_mode
                )
                logger.info(f"Early stopping enabled with patience={config.early_stopping_patience}, "
                           f"min_delta={config.early_stopping_min_delta}, "
                           f"mode={config.early_stopping_mode}")
            
            # Check if training should be resumed from checkpoint
            start_epoch = 1
            if config.resumable_training:
                # Check for latest checkpoint
                checkpoint_files = list(results_manager.checkpoints_dir.glob("checkpoint_epoch_*.pth"))
                if checkpoint_files:
                    # Sort by epoch number
                    checkpoint_files.sort(key=lambda x: int(str(x).split('_')[-1].split('.')[0]))
                    latest_checkpoint = checkpoint_files[-1]
                    
                    logger.info(f"Resuming training from checkpoint: {latest_checkpoint}")
                    checkpoint = torch.load(latest_checkpoint, map_location=device)
                    
                    model.load_state_dict(checkpoint['model_state_dict'])
                    optimizer.load_state_dict(checkpoint['optimizer_state_dict'])
                    
                    # Restore scheduler state if available
                    if scheduler and 'scheduler_state_dict' in checkpoint:
                        scheduler.load_state_dict(checkpoint['scheduler_state_dict'])
                    
                    # Set start epoch
                    start_epoch = checkpoint['epoch'] + 1
                    logger.info(f"Starting from epoch {start_epoch}")
            
            # Set up resource monitor for training
            resource_monitor = ResourceMonitor(log_interval=5)
            resource_monitor.start()
            
            # Training loop
            best_val_accuracy = 0
            best_val_loss = float('inf')
            train_losses = []
            val_losses = []
            val_accuracies = []
            
            # Use timer to measure overall training time
            with TimerContext("Training") as training_timer:
                for epoch in range(start_epoch, config.epochs + 1):
                    # Training phase
                    model.train()
                    running_loss = 0.0
                    
                    # Use timer to measure epoch time
                    with TimerContext(f"Epoch {epoch}") as epoch_timer:
                        for inputs, labels in train_loader:
                            inputs, labels = inputs.to(device), labels.to(device)
                            
                            # Zero gradients
                            optimizer.zero_grad()
                            
                            # Forward pass
                            if model_arch_value == 'siamese':
                                output1, output2 = model(inputs[0], inputs[1])
                                loss = criterion(output1, output2, labels)
                            elif model_arch_value == 'arcface':
                                outputs = model(inputs, labels)
                                loss = criterion(outputs, labels)
                            else:
                                outputs = model(inputs)
                                loss = criterion(outputs, labels)
                            
                            # Backward pass
                            loss.backward()
                            
                            # Apply gradient clipping if enabled
                            if config.use_gradient_clipping:
                                apply_gradient_clipping(
                                    model=model,
                                    max_norm=config.gradient_clipping_max_norm,
                                    adaptive=config.gradient_clipping_adaptive,
                                    model_type=model_arch_value
                                )
                            
                            # Optimize
                            optimizer.step()
                            
                            running_loss += loss.item()
                    
                    epoch_loss = running_loss / len(train_loader)
                    train_losses.append(epoch_loss)
                    
                    # Record training metrics
                    results_manager.record_training_metrics(epoch, {
                        "loss": epoch_loss,
                        "epoch_time": epoch_timer.elapsed_time
                    })
                    
                    # Validation phase
                    model.eval()
                    val_loss = 0.0
                    correct = 0
                    total = 0
                    
                    # For validation calibration (store predictions and ground truth)
                    val_preds = []
                    val_true = []
                    val_scores = []
                    
                    with torch.no_grad():
                        for inputs, labels in val_loader:
                            inputs, labels = inputs.to(device), labels.to(device)
                            
                            if model_arch_value == 'siamese':
                                output1, output2 = model(inputs[0], inputs[1])
                                loss = criterion(output1, output2, labels)
                                preds = (F.pairwise_distance(output1, output2) < 0.5).float()
                                correct += (preds == labels).sum().item()
                                
                                # Store predictions for calibration
                                val_preds.extend(preds.cpu().numpy())
                                val_true.extend(labels.cpu().numpy())
                                val_scores.extend(F.pairwise_distance(output1, output2).cpu().numpy())
                                
                            elif model_arch_value == 'arcface':
                                embeddings = model(inputs)
                                outputs = F.linear(
                                    F.normalize(embeddings), 
                                    F.normalize(model.arcface.weight)
                                )
                                loss = criterion(outputs, labels)
                                _, preds = torch.max(outputs, 1)
                                correct += (preds == labels).sum().item()
                                
                                # Store predictions for calibration
                                val_preds.extend(preds.cpu().numpy())
                                val_true.extend(labels.cpu().numpy())
                                val_scores.extend(F.softmax(outputs, dim=1).cpu().numpy())
                                
                            else:
                                outputs = model(inputs)
                                loss = criterion(outputs, labels)
                                _, preds = torch.max(outputs, 1)
                                correct += (preds == labels).sum().item()
                                
                                # Store predictions for calibration
                                val_preds.extend(preds.cpu().numpy())
                                val_true.extend(labels.cpu().numpy())
                                val_scores.extend(F.softmax(outputs, dim=1).cpu().numpy())
                            
                            val_loss += loss.item()
                            total += labels.size(0)
                    
                    epoch_val_loss = val_loss / len(val_loader)
                    val_losses.append(epoch_val_loss)
                    
                    accuracy = 100 * correct / total
                    val_accuracies.append(accuracy)
                    
                    # Record validation metrics
                    val_metrics = {
                        "loss": epoch_val_loss,
                        "accuracy": accuracy
                    }
                    results_manager.record_evaluation_metrics(epoch, val_metrics)
                    
                    # Record per-class and calibration metrics on validation set periodically
                    if epoch % 5 == 0 or epoch == config.epochs:
                        # Convert predictions array for calibration metrics
                        val_scores_array = np.array(val_scores)
                        
                        # Record validation per-class metrics
                        results_manager.record_per_class_metrics(
                            val_true, val_preds, val_scores_array, class_names, dataset="validation"
                        )
                        
                        # Record validation calibration metrics
                        results_manager.record_calibration_metrics(
                            val_true, val_preds, val_scores_array, dataset="validation"
                        )
                    
                    # Step learning rate scheduler if it's a validation-based scheduler
                    if scheduler:
                        if isinstance(scheduler, torch.optim.lr_scheduler.ReduceLROnPlateau):
                            metric_value = epoch_val_loss if config.early_stopping_metric == "loss" else -accuracy
                            scheduler.step(metric_value)
                    
                    # Save checkpoint if this is the best model (by accuracy or loss)
                    is_best = False
                    if config.early_stopping_mode == 'max':
                        # For metrics like accuracy where higher is better
                        is_best = accuracy > best_val_accuracy
                        if is_best:
                            best_val_accuracy = accuracy
                    else:
                        # For metrics like loss where lower is better
                        is_best = epoch_val_loss < best_val_loss
                        if is_best:
                            best_val_loss = epoch_val_loss
                    
                    # Save checkpoint with appropriate frequency
                    if epoch % config.checkpoint_frequency == 0 or is_best:
                        results_manager.save_model_checkpoint(
                            model=model, 
                            optimizer=optimizer, 
                            epoch=epoch, 
                            is_best=is_best,
                            scheduler=scheduler,
                            metrics=val_metrics
                        )
                    
                    # Log progress
                    logger.info(f'Epoch {epoch}/{config.epochs}, '
                              f'Train Loss: {epoch_loss:.4f}, '
                              f'Val Loss: {epoch_val_loss:.4f}, '
                              f'Accuracy: {accuracy:.2f}%, '
                              f'LR: {optimizer.param_groups[0]["lr"]:.6f}, '
                              f'Time: {epoch_timer.elapsed_time:.2f}s')
                    
                    # Step learning rate scheduler if it's an epoch-based scheduler
                    if scheduler and not isinstance(scheduler, torch.optim.lr_scheduler.ReduceLROnPlateau):
                        scheduler.step()
                    
                    # Check early stopping if enabled
                    if early_stopping:
                        # Use appropriate metric for early stopping
                        if config.early_stopping_metric == "loss":
                            improvement = early_stopping(epoch_val_loss)
                        else:  # Use accuracy
                            # Negate accuracy for 'min' mode
                            es_value = -accuracy if early_stopping.mode == 'min' else accuracy
                            improvement = early_stopping(es_value)
                        
                        # Check if training should be stopped
                        if early_stopping.early_stop:
                            logger.info(f"Early stopping triggered at epoch {epoch}")
                            # Save early stopping trace to a file
                            with open(results_manager.logs_dir / "early_stopping_trace.json", 'w') as f:
                                json.dump({
                                    "trace": early_stopping.trace,
                                    "stopped_epoch": epoch,
                                    "best_score": early_stopping.best_score,
                                    "mode": early_stopping.mode
                                }, f, indent=2)
                            break
            
            # Stop resource monitoring for training
            training_resource_data = resource_monitor.stop()
            training_resource_data["total_time"] = training_timer.elapsed_time
            training_resource_data["epochs_completed"] = epoch - start_epoch + 1
            
            # Record resource metrics for training
            results_manager.record_resource_metrics(training_resource_data, "training")
            
            # Record learning curves
            results_manager.record_learning_curves(train_losses, val_losses, val_accuracies)
            
            # Test the best model
            logger.info("Testing the best model...")
            model.load_state_dict(torch.load(results_manager.checkpoints_dir / "best_model.pth"))
            model.eval()
            
            # Start resource monitoring for inference
            inference_monitor = ResourceMonitor(log_interval=1)
            inference_monitor.start()
            
            # Use timer to measure inference time
            with TimerContext("Inference") as inference_timer:
                # Evaluate on test set
                all_labels = []
                all_preds = []
                all_probs = []
                all_logits = []
                sample_filenames = []  # For tracking which samples are misclassified
                
                with torch.no_grad():
                    for batch_idx, (inputs, labels) in enumerate(test_loader):
                        inputs, labels = inputs.to(device), labels.to(device)
                        
                        # Use timer to measure batch inference time
                        with TimerContext(f"Batch {batch_idx}") as batch_timer:
                            if model_arch_value == 'siamese':
                                output1, output2 = model(inputs[0], inputs[1])
                                distances = F.pairwise_distance(output1, output2)
                                preds = (distances < 0.5).float()
                                
                                # Store predictions
                                all_probs.extend(distances.cpu().numpy())
                                all_logits.extend(distances.cpu().numpy())  # Same as probs for siamese
                                
                            elif model_arch_value == 'arcface':
                                embeddings = model(inputs)
                                logits = F.linear(
                                    F.normalize(embeddings), 
                                    F.normalize(model.arcface.weight)
                                )
                                probs = F.softmax(logits, dim=1)
                                _, preds = torch.max(logits, 1)
                                
                                # Store predictions
                                all_probs.extend(probs.cpu().numpy())
                                all_logits.extend(logits.cpu().numpy())
                                
                            else:
                                logits = model(inputs)
                                probs = F.softmax(logits, dim=1)
                                _, preds = torch.max(logits, 1)
                                
                                # Store predictions
                                all_probs.extend(probs.cpu().numpy())
                                all_logits.extend(logits.cpu().numpy())
                        
                        # Collect batch inference time
                        if batch_idx == 0:
                            # First batch may include warm-up overhead, record separately
                            first_batch_time = batch_timer.elapsed_time
                        
                        all_preds.extend(preds.cpu().numpy())
                        all_labels.extend(labels.cpu().numpy())
                
                # Convert lists to arrays for easier processing
                all_labels = np.array(all_labels)
                all_preds = np.array(all_preds)
                
                # Handle shape of probability outputs based on model type
                if model_arch_value == 'siamese':
                    all_probs = np.array(all_probs).reshape(-1, 1)
                else:
                    all_probs = np.array(all_probs)
                
                all_logits = np.array(all_logits)
            
            # Stop resource monitoring for inference
            inference_resource_data = inference_monitor.stop()
            inference_resource_data["total_time"] = inference_timer.elapsed_time
            inference_resource_data["samples_processed"] = len(all_labels)
            inference_resource_data["samples_per_second"] = len(all_labels) / inference_timer.elapsed_time
            inference_resource_data["first_batch_time"] = first_batch_time
            
            # Record resource metrics for inference
            results_manager.record_resource_metrics(inference_resource_data, "inference")
            
            # Calculate metrics
            from sklearn.metrics import accuracy_score, precision_score, recall_score, f1_score
            
            test_metrics = {
                "accuracy": accuracy_score(all_labels, all_preds) * 100,
                "precision": precision_score(all_labels, all_preds, average='weighted', zero_division=0),
                "recall": recall_score(all_labels, all_preds, average='weighted', zero_division=0),
                "f1_score": f1_score(all_labels, all_preds, average='weighted', zero_division=0),
                "inference_time_total": inference_timer.elapsed_time,
                "inference_time_per_sample": inference_timer.elapsed_time / len(all_labels),
                "inference_samples_per_second": len(all_labels) / inference_timer.elapsed_time
            }
            
            # Record test metrics
            results_manager.record_test_metrics(test_metrics)
            
            # Record confusion matrix with enhanced metrics
            results_manager.record_confusion_matrix(all_labels, all_preds, class_names)
            
            # Record per-class metrics
            results_manager.record_per_class_metrics(all_labels, all_preds, all_probs, class_names)
            
            # Record calibration metrics
            results_manager.record_calibration_metrics(all_labels, all_preds, all_probs)
            
            # Save raw predictions for post-hoc analysis
            results_manager.save_raw_predictions(all_labels, all_preds, all_probs, class_names)
            
            # Include training enhancement details in parameters
            enhancement_params = {
                "preprocessing": config.preprocessing_config.name if config.preprocessing_config else "default",
                "epochs": config.epochs,
                "batch_size": config.batch_size,
                "learning_rate": config.learning_rate,
                "early_stopping": config.use_early_stopping,
                "gradient_clipping": config.use_gradient_clipping,
                "lr_scheduler": config.lr_scheduler_type.value,
                "random_seed": config.random_seed
            }
            
            # Register model with registry
            model_name = f"{model_arch_value}_{config.dataset.value}_{config.experiment_id}"
            self.model_registry.register_model(
                model_name=model_name,
                architecture=model_arch_value,
                dataset_name=config.dataset.value,
                experiment_id=config.experiment_id,
                parameters=enhancement_params,
                metrics=test_metrics,
                checkpoint_path=results_manager.checkpoints_dir / "best_model.pth"
            )
            
            # Generate experiment summary
            summary = results_manager.generate_experiment_summary()
            
            logger.info(f"Experiment completed: {config.experiment_name}")
            logger.info(f"Results saved to: {config.results_dir}")
            
            return summary
        
        except Exception as e:
            logger.error(f"Error running experiment: {str(e)}")
            import traceback
            logger.error(traceback.format_exc())
            results_manager.log_event("experiment_error", {"error": str(e)})
            raise
    
    def _run_cross_dataset_experiment(self, config: ExperimentConfig, 
                                    results_manager: ResultsManager) -> Dict[str, Any]:
        """Run an experiment testing models trained on one dataset against the other."""
        try:
            # Handle model_architecture correctly when it's a list
            model_arch = config.model_architecture
            model_arch_value = model_arch if isinstance(model_arch, list) else model_arch.value
            
            # Create a summary dictionary to store results
            cross_dataset_summary = {
                "experiment_id": config.experiment_id,
                "experiment_name": config.experiment_name,
                "model_architecture": model_arch_value,
                "datasets": ["dataset1", "dataset2"],
                "results": {}
            }
            
            # Run experiments for each dataset
            for dataset in [ExperimentConfig.Dataset.DATASET1, ExperimentConfig.Dataset.DATASET2]:
                # Create a new config for this dataset
                dataset_config = ExperimentConfig(
                    experiment_id=f"{config.experiment_id}_{dataset.value}",
                    experiment_name=f"{config.experiment_name} - {dataset.value}",
                    dataset=dataset,
                    model_architecture=config.model_architecture,
                    preprocessing_config=config.preprocessing_config,
                    epochs=config.epochs,
                    batch_size=config.batch_size,
                    learning_rate=config.learning_rate,
                    cross_dataset_testing=False,
                    results_dir=str(config.results_dir / dataset.value),
                    random_seed=config.random_seed,  # Pass the random seed for reproducibility
                    config_version=config.config_version
                )
                
                # Run experiment for this dataset
                logger.info(f"Running experiment on {dataset.value}...")
                summary = self._run_single_model_experiment(dataset_config, ResultsManager(dataset_config))
                
                # Store results
                cross_dataset_summary["results"][dataset.value] = {
                    "training_summary": summary
                }
                
                # Now test on the other dataset
                other_dataset = ExperimentConfig.Dataset.DATASET1 if dataset == ExperimentConfig.Dataset.DATASET2 else ExperimentConfig.Dataset.DATASET2
                
                logger.info(f"Testing model trained on {dataset.value} against {other_dataset.value}...")
                
                # Get the best model checkpoint
                model_checkpoint = dataset_config.results_dir / "checkpoints" / "best_model.pth"
                
                # Load model
                device = torch.device('cuda' if torch.cuda.is_available() else 'cpu')
                model = get_model(model_arch_value)
                model.load_state_dict(torch.load(model_checkpoint, map_location=device))
                model.to(device)
                model.eval()
                
                # Setup test loader for other dataset with the same random seed
                _, _, test_loader, class_names = self._setup_data_loaders(
                    other_dataset.value, 
                    config.preprocessing_config, 
                    batch_size=config.batch_size,
                    seed=config.random_seed  # Use the same seed for reproducibility
                )
                
                # Test the model
                all_labels = []
                all_preds = []
                
                with torch.no_grad():
                    for inputs, labels in test_loader:
                        inputs, labels = inputs.to(device), labels.to(device)
                        
                        if model_arch_value == 'siamese':
                            output1, output2 = model(inputs[0], inputs[1])
                            preds = (F.pairwise_distance(output1, output2) < 0.5).float()
                        elif model_arch_value == 'arcface':
                            embeddings = model(inputs)
                            outputs = F.linear(
                                F.normalize(embeddings), 
                                F.normalize(model.arcface.weight)
                            )
                            _, preds = torch.max(outputs, 1)
                        else:
                            outputs = model(inputs)
                            _, preds = torch.max(outputs, 1)
                        
                        all_preds.extend(preds.cpu().numpy())
                        all_labels.extend(labels.cpu().numpy())
                
                # Calculate metrics
                from sklearn.metrics import accuracy_score, precision_score, recall_score, f1_score
                
                cross_test_metrics = {
                    "accuracy": accuracy_score(all_labels, all_preds) * 100,
                    "precision": precision_score(all_labels, all_preds, average='weighted', zero_division=0),
                    "recall": recall_score(all_labels, all_preds, average='weighted', zero_division=0),
                    "f1_score": f1_score(all_labels, all_preds, average='weighted', zero_division=0)
                }
                
                # Store cross-dataset results
                cross_dataset_summary["results"][dataset.value]["cross_dataset_testing"] = {
                    "tested_on": other_dataset.value,
                    "metrics": cross_test_metrics
                }
                
                # Log results
                logger.info(f"Cross-dataset testing results for model trained on {dataset.value}:")
                logger.info(f"Accuracy on {other_dataset.value}: {cross_test_metrics['accuracy']:.2f}%")
            
            # Save cross-dataset summary
            cross_summary_path = config.results_dir / "cross_dataset_summary.json"
            with open(cross_summary_path, 'w') as f:
                json.dump(cross_dataset_summary, f, indent=2)
            
            # Log completion
            logger.info(f"Cross-dataset experiment completed: {config.experiment_name}")
            logger.info(f"Results saved to: {config.results_dir}")
            
            return cross_dataset_summary
        
        except Exception as e:
            logger.error(f"Error running cross-dataset experiment: {str(e)}")
            results_manager.log_event("experiment_error", {"error": str(e)})
            raise
    
    def _run_architecture_comparison_experiment(self, config: ExperimentConfig, 
                                             results_manager: ResultsManager) -> Dict[str, Any]:
        """Run an experiment comparing different model architectures."""
        try:
            # Create a summary dictionary to store results
            comparison_summary = {
                "experiment_id": config.experiment_id,
                "experiment_name": config.experiment_name,
                "dataset": config.dataset.value,
                "architectures": [],
                "results": {}
            }
            
            # Determine which architectures to compare
            if isinstance(config.model_architecture, list):
                architectures = config.model_architecture
            elif hasattr(config.model_architecture, 'value'):
                # Handle the case when it's an enum
                architectures = [config.model_architecture.value]
            else:
                # If not specified, compare all architectures
                architectures = [arch.value for arch in ExperimentConfig.ModelArchitecture]
            
            comparison_summary["architectures"] = architectures
            
            # Run experiments for each architecture
            for arch in architectures:
                # Create a new config for this architecture
                arch_config = ExperimentConfig(
                    experiment_id=f"{config.experiment_id}_{arch}",
                    experiment_name=f"{config.experiment_name} - {arch}",
                    dataset=config.dataset,
                    model_architecture=arch,
                    preprocessing_config=config.preprocessing_config,
                    epochs=config.epochs,
                    batch_size=config.batch_size,
                    learning_rate=config.learning_rate,
                    cross_dataset_testing=False,
                    results_dir=str(config.results_dir / arch),
                    random_seed=config.random_seed,  # Pass the random seed for reproducibility
                    config_version=config.config_version
                )
                
                # Run experiment for this architecture
                logger.info(f"Running experiment with {arch} architecture...")
                summary = self._run_single_model_experiment(arch_config, ResultsManager(arch_config))
                
                # Store results
                comparison_summary["results"][arch] = summary
            
            # Create comparison visualization
            self._create_architecture_comparison_plot(comparison_summary, config.results_dir)
            
            # Save comparison summary
            comparison_path = config.results_dir / "architecture_comparison.json"
            with open(comparison_path, 'w') as f:
                json.dump(comparison_summary, f, indent=2)
            
            # Log completion
            logger.info(f"Architecture comparison experiment completed: {config.experiment_name}")
            logger.info(f"Results saved to: {config.results_dir}")
            
            return comparison_summary
        
        except Exception as e:
            logger.error(f"Error running architecture comparison experiment: {str(e)}")
            results_manager.log_event("experiment_error", {"error": str(e)})
            raise
    
    def _create_architecture_comparison_plot(self, comparison_summary: Dict[str, Any], 
                                          output_dir: Path):
        """Create visualization comparing performance of different architectures."""
        architectures = comparison_summary["architectures"]
        
        # Extract metrics
        accuracies = []
        precisions = []
        recalls = []
        f1_scores = []
        
        for arch in architectures:
            if arch in comparison_summary["results"]:
                result = comparison_summary["results"][arch]
                
                # Get test metrics
                if "test_metrics" in result and result["test_metrics"]:
                    metrics = result["test_metrics"][0]  # Get first test metrics entry
                    accuracies.append(metrics.get("accuracy", 0))
                    precisions.append(metrics.get("precision", 0))
                    recalls.append(metrics.get("recall", 0))
                    f1_scores.append(metrics.get("f1_score", 0))
                else:
                    # If no test metrics, add zeros
                    accuracies.append(0)
                    precisions.append(0)
                    recalls.append(0)
                    f1_scores.append(0)
        
        # Create DataFrame
        import pandas as pd
        df = pd.DataFrame({
            "Architecture": architectures,
            "Accuracy": accuracies,
            "Precision": precisions,
            "Recall": recalls,
            "F1 Score": f1_scores
        })
        
        # Save to CSV
        df.to_csv(output_dir / "architecture_comparison.csv", index=False)
        
        # Create bar chart
        plt.figure(figsize=(12, 8))
        
        # Plot metrics
        bar_width = 0.2
        positions = np.arange(len(architectures))
        
        plt.bar(positions - bar_width*1.5, accuracies, bar_width, label='Accuracy', alpha=0.8)
        plt.bar(positions - bar_width/2, precisions, bar_width, label='Precision', alpha=0.8)
        plt.bar(positions + bar_width/2, recalls, bar_width, label='Recall', alpha=0.8)
        plt.bar(positions + bar_width*1.5, f1_scores, bar_width, label='F1 Score', alpha=0.8)
        
        plt.xlabel('Architecture')
        plt.ylabel('Score')
        plt.title('Performance Comparison Across Architectures')
        plt.xticks(positions, architectures)
        plt.legend()
        
        plt.tight_layout()
        plt.savefig(output_dir / "architecture_comparison.png")
        plt.close()
    
    def _setup_data_loaders(self, dataset_name: str, preprocessing_config: PreprocessingConfig, 
                         batch_size: int = 32, seed: int = 42):
        """Set up data loaders for training, validation and testing."""
        from torch.utils.data import DataLoader
        from torchvision import datasets, transforms
        import torch.utils.data
        
        # Set generator with specific seed for reproducible data splitting
        g = torch.Generator()
        g.manual_seed(seed)
        
        # Define transforms
        transform = transforms.Compose([
            transforms.Resize((224, 224)),
            transforms.ToTensor(),
            transforms.Normalize(mean=[0.485, 0.456, 0.406],
                               std=[0.229, 0.224, 0.225])
        ])
        
        # If preprocessing config is provided, use the corresponding processed data
        if preprocessing_config:
            data_dir = Path(f"processed_data/{preprocessing_config.name}/{dataset_name}")
        else:
            # Otherwise, use default processed data
            data_dir = Path(f"processed_data/default/{dataset_name}")
        
        # Check if directory exists
        if not data_dir.exists():
            raise ValueError(f"Data directory not found: {data_dir}")
        
        # Load datasets
        train_dataset = datasets.ImageFolder(data_dir / "train", transform=transform)
        val_dataset = datasets.ImageFolder(data_dir / "val", transform=transform)
        test_dataset = datasets.ImageFolder(data_dir / "test", transform=transform)
        
        # Create data loaders with seed-controlled randomness
        train_loader = DataLoader(
            train_dataset, 
            batch_size=batch_size, 
            shuffle=True, 
            num_workers=4, 
            generator=g,
            worker_init_fn=lambda worker_id: torch.initial_seed() + worker_id
        )
        val_loader = DataLoader(
            val_dataset, 
            batch_size=batch_size, 
            shuffle=False, 
            num_workers=4
        )
        test_loader = DataLoader(
            test_dataset, 
            batch_size=batch_size, 
            shuffle=False, 
            num_workers=4
        )
        
        return train_loader, val_loader, test_loader, train_dataset.classes 


class DatasetComparisonExperiment:
    """Experiment for comparing performance between datasets."""
    
    def __init__(self, experiment_manager: ExperimentManager, 
                model_architecture: str = "cnn", 
                preprocessing_config: Optional[PreprocessingConfig] = None):
        """Initialize dataset comparison experiment."""
        self.experiment_manager = experiment_manager
        self.model_architecture = model_architecture
        self.preprocessing_config = preprocessing_config
        
    def run(self, epochs: int = 30, batch_size: int = 32, 
           learning_rate: float = 0.001) -> Dict[str, Any]:
        """Run dataset comparison experiment."""
        # Create experiment configuration
        config = ExperimentConfig(
            experiment_name=f"Dataset Comparison - {self.model_architecture}",
            dataset=ExperimentConfig.Dataset.BOTH,
            model_architecture=self.model_architecture,
            preprocessing_config=self.preprocessing_config,
            epochs=epochs,
            batch_size=batch_size,
            learning_rate=learning_rate,
            cross_dataset_testing=True
        )
        
        # Run experiment
        return self.experiment_manager.run_experiment(config)
    
    def analyze_results(self, results_dir: Path) -> Dict[str, Any]:
        """Analyze and visualize dataset comparison results."""
        # Load experiment results
        try:
            with open(results_dir / "cross_dataset_summary.json", 'r') as f:
                summary = json.load(f)
        except FileNotFoundError:
            raise ValueError(f"Cross-dataset summary not found in {results_dir}")
        
        # Extract metrics for comparison
        dataset1_metrics = summary["results"]["dataset1"]["training_summary"]["test_metrics"][0]
        dataset2_metrics = summary["results"]["dataset2"]["training_summary"]["test_metrics"][0]
        dataset1_on_2_metrics = summary["results"]["dataset1"]["cross_dataset_testing"]["metrics"]
        dataset2_on_1_metrics = summary["results"]["dataset2"]["cross_dataset_testing"]["metrics"]
        
        # Create comparison table
        comparison = {
            "dataset1_same_domain": dataset1_metrics,
            "dataset2_same_domain": dataset2_metrics,
            "dataset1_cross_domain": dataset1_on_2_metrics,
            "dataset2_cross_domain": dataset2_on_1_metrics
        }
        
        # Create visualization
        self._create_dataset_comparison_plot(comparison, results_dir)
        
        # Analyze dataset statistics
        dataset_stats = self._analyze_dataset_statistics()
        
        # Combine results
        analysis = {
            "comparison": comparison,
            "dataset_statistics": dataset_stats,
            "insights": self._generate_insights(comparison, dataset_stats)
        }
        
        # Save analysis
        with open(results_dir / "dataset_comparison_analysis.json", 'w') as f:
            json.dump(analysis, f, indent=2)
        
        return analysis
    
    def _create_dataset_comparison_plot(self, comparison: Dict[str, Dict[str, float]], 
                                     results_dir: Path):
        """Create visualization comparing performance across datasets."""
        # Extract accuracy metrics
        accuracies = [
            comparison["dataset1_same_domain"]["accuracy"],
            comparison["dataset1_cross_domain"]["accuracy"],
            comparison["dataset2_same_domain"]["accuracy"],
            comparison["dataset2_cross_domain"]["accuracy"]
        ]
        
        # Create bar chart
        plt.figure(figsize=(10, 6))
        
        bar_labels = [
            "Dataset 1\nSame Domain",
            "Dataset 1\nCross Domain",
            "Dataset 2\nSame Domain",
            "Dataset 2\nCross Domain"
        ]
        
        bars = plt.bar(bar_labels, accuracies, color=['blue', 'lightblue', 'red', 'lightcoral'])
        
        plt.ylabel('Accuracy (%)')
        plt.title('Cross-Dataset Performance Comparison')
        
        # Add values on top of bars
        for bar in bars:
            height = bar.get_height()
            plt.annotate(f'{height:.1f}%',
                       xy=(bar.get_x() + bar.get_width() / 2, height),
                       xytext=(0, 3),
                       textcoords="offset points",
                       ha='center', va='bottom')
        
        plt.tight_layout()
        plt.savefig(results_dir / "dataset_comparison.png")
        plt.close()
        
        # Create a more detailed metrics comparison
        metrics = ['accuracy', 'precision', 'recall', 'f1_score']
        metric_names = ['Accuracy', 'Precision', 'Recall', 'F1 Score']
        
        plt.figure(figsize=(12, 8))
        
        x = np.arange(len(metrics))
        width = 0.2
        
        plt.bar(x - width*1.5, [comparison["dataset1_same_domain"][m] for m in metrics], 
              width, label='Dataset 1 Same Domain', color='blue')
        plt.bar(x - width/2, [comparison["dataset1_cross_domain"][m] for m in metrics], 
              width, label='Dataset 1 Cross Domain', color='lightblue')
        plt.bar(x + width/2, [comparison["dataset2_same_domain"][m] for m in metrics], 
              width, label='Dataset 2 Same Domain', color='red')
        plt.bar(x + width*1.5, [comparison["dataset2_cross_domain"][m] for m in metrics], 
              width, label='Dataset 2 Cross Domain', color='lightcoral')
        
        plt.xlabel('Metric')
        plt.ylabel('Score')
        plt.title('Detailed Cross-Dataset Performance Comparison')
        plt.xticks(x, metric_names)
        plt.legend()
        
        plt.tight_layout()
        plt.savefig(results_dir / "dataset_comparison_detailed.png")
        plt.close()
    
    def _analyze_dataset_statistics(self) -> Dict[str, Any]:
        """Analyze and compare dataset statistics."""
        # This could include:
        # - Number of images per dataset
        # - Number of subjects
        # - Images per subject
        # - Image quality metrics
        # - Face detection success rates
        # - etc.
        
        # For the sake of this implementation, we'll return some dummy data
        # In a real implementation, you would analyze the actual datasets
        return {
            "dataset1": {
                "num_subjects": 36,
                "images_per_subject": 49,
                "total_images": 36 * 49,
                "image_quality": "varied"
            },
            "dataset2": {
                "num_subjects": 18,
                "images_per_subject": 100,
                "total_images": 18 * 100,
                "image_quality": "consistent"
            }
        }
    
    def _generate_insights(self, comparison: Dict[str, Dict[str, float]], 
                         dataset_stats: Dict[str, Dict[str, Any]]) -> List[str]:
        """Generate insights from dataset comparison."""
        insights = []
        
        # Compare same-domain performance
        dataset1_accuracy = comparison["dataset1_same_domain"]["accuracy"]
        dataset2_accuracy = comparison["dataset2_same_domain"]["accuracy"]
        
        if dataset1_accuracy > dataset2_accuracy:
            insights.append(f"Dataset 1 (high diversity) shows better in-domain performance than Dataset 2 (high quantity) by {dataset1_accuracy - dataset2_accuracy:.1f}%.")
        else:
            insights.append(f"Dataset 2 (high quantity) shows better in-domain performance than Dataset 1 (high diversity) by {dataset2_accuracy - dataset1_accuracy:.1f}%.")
        
        # Compare cross-domain generalization
        dataset1_cross = comparison["dataset1_cross_domain"]["accuracy"]
        dataset2_cross = comparison["dataset2_cross_domain"]["accuracy"]
        
        if dataset1_cross > dataset2_cross:
            insights.append(f"Models trained on Dataset 1 (high diversity) generalize better to new datasets by {dataset1_cross - dataset2_cross:.1f}%.")
        else:
            insights.append(f"Models trained on Dataset 2 (high quantity) generalize better to new datasets by {dataset2_cross - dataset1_cross:.1f}%.")
        
        # Compare performance drop
        drop1 = dataset1_accuracy - dataset1_cross
        drop2 = dataset2_accuracy - dataset2_cross
        
        if drop1 < drop2:
            insights.append(f"Models trained on Dataset 1 (high diversity) show more robust cross-domain performance with {drop2 - drop1:.1f}% less performance drop.")
        else:
            insights.append(f"Models trained on Dataset 2 (high quantity) show more robust cross-domain performance with {drop1 - drop2:.1f}% less performance drop.")
        
        return insights


class CrossArchitectureExperiment:
    """Experiment for comparing different model architectures."""
    
    def __init__(self, experiment_manager: ExperimentManager, 
                architectures: Optional[List[str]] = None,
                dataset: str = "both",
                preprocessing_config: Optional[PreprocessingConfig] = None):
        """Initialize architecture comparison experiment."""
        self.experiment_manager = experiment_manager
        
        # If architectures not specified, compare all available architectures
        if architectures is None:
            self.architectures = [arch.value for arch in ExperimentConfig.ModelArchitecture]
        else:
            self.architectures = architectures
            
        self.dataset = dataset
        self.preprocessing_config = preprocessing_config
        
    def run(self, epochs: int = 30, batch_size: int = 32, 
           learning_rate: float = 0.001) -> Dict[str, Any]:
        """Run architecture comparison experiment."""
        # Create experiment configuration
        config = ExperimentConfig(
            experiment_name=f"Architecture Comparison - {self.dataset}",
            dataset=self.dataset,
            model_architecture=self.architectures,  # Pass all architectures to test
            preprocessing_config=self.preprocessing_config,
            epochs=epochs,
            batch_size=batch_size,
            learning_rate=learning_rate,
            cross_dataset_testing=False,
            results_dir=str(self.results_dir / "architecture_comparison")
        )
        
        # Run experiment
        return self.experiment_manager.run_experiment(config)
    
    def analyze_results(self, results_dir: Path) -> Dict[str, Any]:
        """Analyze and visualize architecture comparison results."""
        # Load experiment results
        try:
            with open(results_dir / "architecture_comparison.json", 'r') as f:
                summary = json.load(f)
        except FileNotFoundError:
            raise ValueError(f"Architecture comparison summary not found in {results_dir}")
        
        # Extract metrics for each architecture
        architectures = summary["architectures"]
        metrics = {}
        
        for arch in architectures:
            if arch in summary["results"]:
                result = summary["results"][arch]
                if "test_metrics" in result and result["test_metrics"]:
                    metrics[arch] = result["test_metrics"][0]
        
        # Create comprehensive comparison visualization
        self._create_architecture_comparison_plot(metrics, results_dir)
        
        # Create comparative performance matrix
        perf_matrix = self._create_performance_matrix(metrics)
        
        # Generate architecture insights
        insights = self._generate_architecture_insights(metrics)
        
        # Combine results
        analysis = {
            "metrics": metrics,
            "performance_matrix": perf_matrix,
            "insights": insights
        }
        
        # Save analysis
        with open(results_dir / "architecture_analysis.json", 'w') as f:
            json.dump(analysis, f, indent=2)
        
        return analysis
    
    def _create_architecture_comparison_plot(self, metrics: Dict[str, Dict[str, float]], 
                                         results_dir: Path):
        """Create detailed visualization comparing architectures."""
        # Extract metric types
        metric_types = ['accuracy', 'precision', 'recall', 'f1_score']
        architectures = list(metrics.keys())
        
        # Create a figure with subplots
        fig, axes = plt.subplots(2, 2, figsize=(16, 12))
        axes = axes.flatten()
        
        # Plot each metric type
        for i, metric in enumerate(metric_types):
            values = [metrics[arch].get(metric, 0) for arch in architectures]
            
            axes[i].bar(architectures, values, color='skyblue')
            axes[i].set_title(f'{metric.capitalize()} Comparison')
            axes[i].set_ylabel('Score')
            axes[i].set_xticklabels(architectures, rotation=45)
            
            # Add values on top of bars
            for j, v in enumerate(values):
                axes[i].text(j, v + 0.01, f'{v:.3f}', ha='center')
        
        plt.tight_layout()
        plt.savefig(results_dir / "architecture_comparison_detailed.png")
        plt.close()
        
        # Create a spider/radar chart for comparative visualization
        categories = metric_types
        fig = plt.figure(figsize=(10, 10))
        
        # Create radar chart
        ax = fig.add_subplot(111, polar=True)
        
        # Compute angle for each category
        angles = np.linspace(0, 2*np.pi, len(categories), endpoint=False).tolist()
        angles += angles[:1]  # Close the loop
        
        # Plot each architecture
        for arch in architectures:
            values = [metrics[arch].get(metric, 0) for metric in categories]
            values += values[:1]  # Close the loop
            
            ax.plot(angles, values, linewidth=2, label=arch)
            ax.fill(angles, values, alpha=0.1)
        
        # Set category labels
        plt.xticks(angles[:-1], categories)
        
        # Add legend
        plt.legend(loc='upper right')
        
        plt.title('Architecture Comparison - Radar Chart')
        plt.savefig(results_dir / "architecture_comparison_radar.png")
        plt.close()
    
    def _create_performance_matrix(self, metrics: Dict[str, Dict[str, float]]) -> Dict[str, Dict[str, float]]:
        """Create a performance matrix comparing architectures on key metrics."""
        architectures = list(metrics.keys())
        matrix = {}
        
        # Create a matrix with relative performance rankings
        for metric in ['accuracy', 'precision', 'recall', 'f1_score']:
            # Sort architectures by this metric
            sorted_archs = sorted(architectures, 
                                key=lambda arch: metrics[arch].get(metric, 0), 
                                reverse=True)
            
            # Assign rankings
            for i, arch in enumerate(sorted_archs):
                if arch not in matrix:
                    matrix[arch] = {}
                
                matrix[arch][metric] = i + 1  # 1-based ranking
        
        # Calculate overall ranking
        for arch in architectures:
            if arch in matrix:
                matrix[arch]['overall'] = sum([matrix[arch].get(m, len(architectures)) 
                                            for m in ['accuracy', 'precision', 'recall', 'f1_score']])
        
        return matrix
    
    def _generate_architecture_insights(self, metrics: Dict[str, Dict[str, float]]) -> List[str]:
        """Generate insights from architecture comparison."""
        insights = []
        
        # Find best overall architecture
        if metrics:
            best_arch = max(metrics.keys(), key=lambda arch: metrics[arch].get('accuracy', 0))
            best_acc = metrics[best_arch].get('accuracy', 0)
            insights.append(f"{best_arch} achieves the highest accuracy at {best_acc:.2f}%.")
        
        # Compare architectures pairwise
        architectures = list(metrics.keys())
        for i, arch1 in enumerate(architectures):
            for arch2 in architectures[i+1:]:
                acc1 = metrics[arch1].get('accuracy', 0)
                acc2 = metrics[arch2].get('accuracy', 0)
                
                if abs(acc1 - acc2) > 1.0:  # Only report meaningful differences
                    better = arch1 if acc1 > acc2 else arch2
                    worse = arch2 if acc1 > acc2 else arch1
                    diff = abs(acc1 - acc2)
                    insights.append(f"{better} outperforms {worse} by {diff:.2f}% in accuracy.")
        
        # Identify strengths for each architecture
        for arch in architectures:
            # Find the metric where this architecture performs best relative to others
            best_metric = None
            best_rank = float('inf')
            
            for metric in ['accuracy', 'precision', 'recall', 'f1_score']:
                # Calculate the rank of this architecture for this metric
                values = [metrics[a].get(metric, 0) for a in architectures]
                sorted_values = sorted(values, reverse=True)
                rank = sorted_values.index(metrics[arch].get(metric, 0)) + 1
                
                if rank < best_rank:
                    best_rank = rank
                    best_metric = metric
            
            if best_metric and best_rank <= 2:  # Only highlight top performers
                insights.append(f"{arch} shows particular strength in {best_metric} ({metrics[arch].get(best_metric, 0):.2f}).")
        
        return insights


class HyperparameterExperiment:
    """Experiment for optimizing hyperparameters."""
    
    def __init__(self, experiment_manager: ExperimentManager,
                model_architecture: str = "cnn",
                dataset: str = "both",
                preprocessing_config: Optional[PreprocessingConfig] = None,
                n_trials: int = 20,
                timeout: int = 7200):  # 2 hours default timeout
        """Initialize hyperparameter experiment."""
        self.experiment_manager = experiment_manager
        self.model_architecture = model_architecture
        self.dataset = dataset
        self.preprocessing_config = preprocessing_config
        self.n_trials = n_trials
        self.timeout = timeout
        
    def run(self) -> Dict[str, Any]:
        """Run hyperparameter optimization experiment."""
        # Create a timestamp for this optimization run
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        experiment_id = f"hyperopt_{self.model_architecture}_{timestamp}"
        
        # Create results directory
        results_dir = OUT_DIR / experiment_id
        results_dir.mkdir(parents=True, exist_ok=True)
        
        # Create study
        study = optuna.create_study(
            direction="maximize",
            study_name=f"hyperopt_{self.model_architecture}",
            storage=f"sqlite:///{results_dir}/optuna.db",
            load_if_exists=True
        )
        
        # Run optimization
        study.optimize(
            lambda trial: self._objective(trial, results_dir),
            n_trials=self.n_trials,
            timeout=self.timeout
        )
        
        # Save study results
        self._save_study_results(study, results_dir)
        
        # Get best parameters
        best_params = study.best_params
        
        # Run a final experiment with best parameters
        best_config = ExperimentConfig(
            experiment_id=f"{experiment_id}_best",
            experiment_name=f"Best Hyperparameters for {self.model_architecture}",
            dataset=self.dataset,
            model_architecture=self.model_architecture,
            preprocessing_config=self.preprocessing_config,
            epochs=best_params.get("epochs", 30),
            batch_size=best_params.get("batch_size", 32),
            learning_rate=best_params.get("learning_rate", 0.001),
            cross_dataset_testing=False,
            results_dir=str(results_dir / "best_params")
        )
        
        # Run experiment with best parameters
        final_results = self.experiment_manager.run_experiment(best_config)
        
        # Combine results
        summary = {
            "study_name": study.study_name,
            "num_trials": len(study.trials),
            "best_trial": study.best_trial.number,
            "best_params": best_params,
            "best_value": study.best_value,
            "final_results": final_results
        }
        
        # Save summary
        with open(results_dir / "hyperopt_summary.json", 'w') as f:
            json.dump(summary, f, indent=2)
        
        return summary
    
    def _objective(self, trial: optuna.Trial, results_dir: Path) -> float:
        """Objective function for hyperparameter optimization."""
        # Define hyperparameter space
        epochs = trial.suggest_int("epochs", 10, 50)
        batch_size = trial.suggest_categorical("batch_size", [16, 32, 64, 128])
        learning_rate = trial.suggest_float("learning_rate", 1e-5, 1e-2, log=True)
        
        # Additional hyperparameters can be added here
        
        # Create experiment configuration
        config = ExperimentConfig(
            experiment_id=f"trial_{trial.number}",
            experiment_name=f"Hyperopt Trial {trial.number}",
            dataset=self.dataset,
            model_architecture=self.model_architecture,
            preprocessing_config=self.preprocessing_config,
            epochs=epochs,
            batch_size=batch_size,
            learning_rate=learning_rate,
            cross_dataset_testing=False,
            results_dir=str(results_dir / f"trial_{trial.number}")
        )
        
        try:
            # Run experiment
            results = self.experiment_manager.run_experiment(config)
            
            # Extract validation accuracy from results
            if "best_validation_metrics" in results and "accuracy" in results["best_validation_metrics"]:
                return results["best_validation_metrics"]["accuracy"]
            else:
                # If validation metrics not found, use test accuracy
                if "test_metrics" in results and results["test_metrics"] and "accuracy" in results["test_metrics"][0]:
                    return results["test_metrics"][0]["accuracy"]
                else:
                    # If no metrics found, return a very low score
                    return 0.0
        
        except Exception as e:
            logger.error(f"Error in hyperopt trial {trial.number}: {str(e)}")
            return 0.0  # Return low score on error
    
    def _save_study_results(self, study: optuna.Study, results_dir: Path):
        """Save and visualize optimization study results."""
        # Save optimization history
        plt.figure(figsize=(10, 6))
        optuna.visualization.matplotlib.plot_optimization_history(study)
        plt.savefig(results_dir / "optimization_history.png")
        plt.close()
        
        # Save parameter importances
        plt.figure(figsize=(10, 6))
        try:
            optuna.visualization.matplotlib.plot_param_importances(study)
            plt.savefig(results_dir / "param_importances.png")
        except:
            logger.warning("Could not plot parameter importances (requires more than one trial)")
        plt.close()
        
        # Save parallel coordinate plot
        plt.figure(figsize=(12, 8))
        try:
            optuna.visualization.matplotlib.plot_parallel_coordinate(study)
            plt.savefig(results_dir / "parallel_coordinate.png")
        except:
            logger.warning("Could not plot parallel coordinates (requires more than one trial)")
        plt.close()
        
        # Create optimization trials table
        trials_data = []
        for trial in study.trials:
            if trial.state == optuna.trial.TrialState.COMPLETE:
                trial_data = {
                    "number": trial.number,
                    "value": trial.value,
                    **trial.params
                }
                trials_data.append(trial_data)
        
        # Save trials data to CSV
        if trials_data:
            pd.DataFrame(trials_data).to_csv(results_dir / "trials.csv", index=False)


class ResultsCompiler:
    """Utility for creating presentation-ready reports and visualizations."""
    
    def __init__(self, base_results_dir: Optional[Path] = None):
        """Initialize results compiler."""
        self.base_results_dir = base_results_dir or OUT_DIR
    
    def compile_experiment_summary(self, experiment_id: str) -> Dict[str, Any]:
        """Compile a comprehensive summary of an experiment."""
        # Find experiment directory
        experiment_dirs = list(self.base_results_dir.glob(f"*{experiment_id}*"))
        if not experiment_dirs:
            raise ValueError(f"No experiment found with ID: {experiment_id}")
        
        experiment_dir = experiment_dirs[0]
        
        # Load experiment configuration
        config_file = experiment_dir / "experiment_config.json"
        if not config_file.exists():
            config_file = next(experiment_dir.glob("**/experiment_config.json"), None)
            
        if not config_file:
            raise ValueError(f"No configuration file found for experiment: {experiment_id}")
        
        with open(config_file, 'r') as f:
            config = json.load(f)
        
        # Load experiment summary
        summary_file = experiment_dir / "experiment_summary.json"
        if not summary_file.exists():
            summary_file = next(experiment_dir.glob("**/experiment_summary.json"), None)
            
        summary = None
        if summary_file:
            with open(summary_file, 'r') as f:
                summary = json.load(f)
        
        # Find all visualization files
        vis_files = list(experiment_dir.glob("**/*.png"))
        
        # Create report structure
        report = {
            "experiment_id": experiment_id,
            "config": config,
            "summary": summary,
            "visualizations": [str(f.relative_to(self.base_results_dir)) for f in vis_files]
        }
        
        return report
    
    def generate_comparative_report(self, experiment_ids: List[str], 
                                   output_dir: Optional[Path] = None) -> Path:
        """Generate a comparative report of multiple experiments."""
        if output_dir is None:
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            output_dir = self.base_results_dir / f"comparative_report_{timestamp}"
            
        output_dir.mkdir(parents=True, exist_ok=True)
        
        # Compile summaries for all experiments
        experiment_summaries = []
        for exp_id in experiment_ids:
            try:
                summary = self.compile_experiment_summary(exp_id)
                experiment_summaries.append(summary)
            except ValueError as e:
                logger.warning(f"Error compiling summary for experiment {exp_id}: {str(e)}")
        
        # Create comparative visualizations
        self._create_comparative_visualizations(experiment_summaries, output_dir)
        
        # Create report content
        report = {
            "generated_at": datetime.datetime.now().isoformat(),
            "experiments": experiment_summaries,
            "comparative_visualizations": list(output_dir.glob("*.png"))
        }
        
        # Save report
        with open(output_dir / "comparative_report.json", 'w') as f:
            json.dump(report, f, indent=2)
        
        return output_dir
    
    def _create_comparative_visualizations(self, summaries: List[Dict[str, Any]], 
                                        output_dir: Path):
        """Create visualizations comparing multiple experiments."""
        # Extract experiment names and key metrics
        exp_names = []
        accuracies = []
        precisions = []
        recalls = []
        f1_scores = []
        
        for summary in summaries:
            if "summary" in summary and summary["summary"] and "test_metrics" in summary["summary"]:
                test_metrics = summary["summary"]["test_metrics"]
                if test_metrics and len(test_metrics) > 0:
                    metrics = test_metrics[0]
                    
                    exp_names.append(summary["experiment_id"])
                    accuracies.append(metrics.get("accuracy", 0))
                    precisions.append(metrics.get("precision", 0))
                    recalls.append(metrics.get("recall", 0))
                    f1_scores.append(metrics.get("f1_score", 0))
        
        if not exp_names:
            logger.warning("No valid metrics found for comparative visualization")
            return
        
        # Create comparison bar chart
        plt.figure(figsize=(12, 8))
        
        x = np.arange(len(exp_names))
        width = 0.2
        
        plt.bar(x - width*1.5, accuracies, width, label='Accuracy')
        plt.bar(x - width/2, precisions, width, label='Precision')
        plt.bar(x + width/2, recalls, width, label='Recall')
        plt.bar(x + width*1.5, f1_scores, width, label='F1 Score')
        
        plt.xlabel('Experiment')
        plt.ylabel('Score')
        plt.title('Performance Comparison Across Experiments')
        plt.xticks(x, exp_names, rotation=45)
        plt.legend()
        
        plt.tight_layout()
        plt.savefig(output_dir / "performance_comparison.png")
        plt.close()
    
    def generate_powerpoint_report(self, experiment_ids: List[str], 
                                 output_path: Optional[Path] = None) -> Path:
        """Generate a PowerPoint presentation with experiment results."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            raise ImportError("python-pptx is required for PowerPoint generation. Install with: pip install python-pptx")
        
        if output_path is None:
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = self.base_results_dir / f"experiment_report_{timestamp}.pptx"
        
        # Create a presentation
        prs = Presentation()
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Face Recognition Experiment Results"
        subtitle.text = f"Generated on {datetime.datetime.now().strftime('%Y-%m-%d %H:%M')}"
        
        # Compile summaries for all experiments
        for exp_id in experiment_ids:
            try:
                summary = self.compile_experiment_summary(exp_id)
                
                # Add experiment slide
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title = slide.shapes.title
                title.text = f"Experiment: {summary['experiment_id']}"
                
                # Add experiment details
                if "config" in summary:
                    config = summary["config"]
                    content = slide.placeholders[1]
                    tf = content.text_frame
                    tf.text = f"Model: {config.get('model_architecture', 'N/A')}\n"
                    tf.text += f"Dataset: {config.get('dataset', 'N/A')}\n"
                    
                    if "summary" in summary and summary["summary"]:
                        s = summary["summary"]
                        if "test_metrics" in s and s["test_metrics"] and len(s["test_metrics"]) > 0:
                            metrics = s["test_metrics"][0]
                            tf.text += f"\nResults:\n"
                            tf.text += f"Accuracy: {metrics.get('accuracy', 0):.2f}%\n"
                            tf.text += f"Precision: {metrics.get('precision', 0):.2f}\n"
                            tf.text += f"Recall: {metrics.get('recall', 0):.2f}\n"
                            tf.text += f"F1 Score: {metrics.get('f1_score', 0):.2f}\n"
                
                # Add visualization slides
                for vis in summary.get("visualizations", [])[:3]:  # Limit to 3 visualizations
                    vis_path = self.base_results_dir / vis
                    if vis_path.exists():
                        slide = prs.slides.add_slide(prs.slide_layouts[6])
                        title = slide.shapes.title
                        title.text = f"Visualization: {vis_path.stem}"
                        
                        # Add image
                        slide.shapes.add_picture(str(vis_path), Inches(1), Inches(1.5), 
                                             width=Inches(8), height=Inches(5))
            
            except ValueError as e:
                logger.warning(f"Error compiling summary for experiment {exp_id}: {str(e)}")
                continue
        
        # Save presentation
        prs.save(output_path)
        
        return output_path
    
    def export_to_excel(self, experiment_ids: List[str], 
                      output_path: Optional[Path] = None) -> Path:
        """Export experiment results to Excel spreadsheet."""
        if output_path is None:
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = self.base_results_dir / f"experiment_results_{timestamp}.xlsx"
        
        # Compile results for all experiments
        results = []
        
        for exp_id in experiment_ids:
            try:
                summary = self.compile_experiment_summary(exp_id)
                
                if "config" in summary and "summary" in summary and summary["summary"]:
                    config = summary["config"]
                    s = summary["summary"]
                    
                    result = {
                        "experiment_id": exp_id,
                        "model_architecture": config.get("model_architecture", ""),
                        "dataset": config.get("dataset", ""),
                        "epochs": config.get("epochs", ""),
                        "batch_size": config.get("batch_size", ""),
                        "learning_rate": config.get("learning_rate", "")
                    }
                    
                    # Add test metrics
                    if "test_metrics" in s and s["test_metrics"] and len(s["test_metrics"]) > 0:
                        metrics = s["test_metrics"][0]
                        result.update({
                            "accuracy": metrics.get("accuracy", ""),
                            "precision": metrics.get("precision", ""),
                            "recall": metrics.get("recall", ""),
                            "f1_score": metrics.get("f1_score", "")
                        })
                    
                    results.append(result)
            
            except ValueError as e:
                logger.warning(f"Error compiling summary for experiment {exp_id}: {str(e)}")
                continue
        
        # Create Excel file
        if results:
            df = pd.DataFrame(results)
            df.to_excel(output_path, index=False)
        
        return output_path


def create_command_line_interface():
    """Create a command-line interface for the experiment manager."""
    import argparse
    
    parser = argparse.ArgumentParser(description="Face Recognition Experiment Manager")
    subparsers = parser.add_subparsers(dest="command", help="Command to run")
    
    # Run experiment from config file
    parser_run = subparsers.add_parser("run", 
                                     help="Run experiment from config file")
    parser_run.add_argument("config_file", type=str,
                          help="Path to experiment configuration file (JSON or YAML)")
    
    # Dataset Comparison experiment
    parser_dataset = subparsers.add_parser("dataset-comparison", 
                                          help="Run dataset comparison experiment")
    parser_dataset.add_argument("--architecture", type=str, default="cnn",
                              help="Model architecture to use")
    parser_dataset.add_argument("--epochs", type=int, default=30,
                              help="Number of training epochs")
    parser_dataset.add_argument("--batch-size", type=int, default=32,
                              help="Batch size for training")
    parser_dataset.add_argument("--learning-rate", type=float, default=0.001,
                              help="Learning rate for training")
    parser_dataset.add_argument("--output-config", type=str,
                              help="Save configuration to file (specify .json or .yaml extension)")
    
    # Architecture Comparison experiment
    parser_arch = subparsers.add_parser("architecture-comparison", 
                                      help="Run architecture comparison experiment")
    parser_arch.add_argument("--dataset", type=str, default="both",
                           help="Dataset to use (dataset1, dataset2, or both)")
    parser_arch.add_argument("--architectures", type=str, nargs="+",
                           default=["baseline", "cnn", "attention", "arcface", "hybrid"],
                           help="Model architectures to compare")
    parser_arch.add_argument("--epochs", type=int, default=30,
                           help="Number of training epochs")
    parser_arch.add_argument("--batch-size", type=int, default=32,
                           help="Batch size for training")
    parser_arch.add_argument("--learning-rate", type=float, default=0.001,
                           help="Learning rate for training")
    parser_arch.add_argument("--output-config", type=str,
                           help="Save configuration to file (specify .json or .yaml extension)")
    
    # Hyperparameter Optimization experiment
    parser_hyperopt = subparsers.add_parser("hyperparameter-optimization",
                                          help="Run hyperparameter optimization experiment")
    parser_hyperopt.add_argument("--architecture", type=str, default="cnn",
                               help="Model architecture to optimize")
    parser_hyperopt.add_argument("--dataset", type=str, default="both",
                               help="Dataset to use (dataset1, dataset2, or both)")
    parser_hyperopt.add_argument("--n-trials", type=int, default=20,
                               help="Number of optimization trials")
    parser_hyperopt.add_argument("--timeout", type=int, default=7200,
                               help="Timeout in seconds")
    parser_hyperopt.add_argument("--output-config", type=str,
                               help="Save configuration to file (specify .json or .yaml extension)")
    
    # Generate report
    parser_report = subparsers.add_parser("generate-report",
                                        help="Generate report from experiment results")
    parser_report.add_argument("--experiment-ids", type=str, nargs="+", required=True,
                             help="Experiment IDs to include in report")
    parser_report.add_argument("--format", type=str, choices=["json", "pptx", "excel", "yaml"],
                             default="json", help="Report format")
    parser_report.add_argument("--output", type=str, help="Output file path")
    
    # Resume experiment
    parser_resume = subparsers.add_parser("resume",
                                        help="Resume an interrupted experiment")
    parser_resume.add_argument("experiment_id", type=str,
                             help="ID of experiment to resume")
    
    return parser


def main():
    parser = create_command_line_interface()
    args = parser.parse_args()
    
    # Create experiment manager
    experiment_manager = ExperimentManager()

    if args.command == "run":
        # Run experiment from config file
        results = experiment_manager.run_experiment(args.config_file)
        print(f"Experiment Results: {results}")
    elif args.command == "dataset-comparison":
        # Create dataset comparison experiment
        experiment = DatasetComparisonExperiment(
            experiment_manager=experiment_manager,
            model_architecture=args.architecture,
        )
        
        # If output config is specified, save the configuration
        if hasattr(args, 'output_config') and args.output_config:
            # Create config
            config = ExperimentConfig(
                experiment_name=f"Dataset Comparison - {args.architecture}",
                dataset=ExperimentConfig.Dataset.BOTH,
                model_architecture=args.architecture,
                epochs=args.epochs,
                batch_size=args.batch_size,
                learning_rate=args.learning_rate,
                cross_dataset_testing=True
            )
            
            # Save config file
            output_path = Path(args.output_config)
            if output_path.suffix.lower() in ['.yaml', '.yml']:
                config.save_yaml(output_path)
            else:
                config.save(output_path)
            print(f"Configuration saved to: {output_path}")
            return
        
        # Run experiment
        results = experiment.run(
            epochs=args.epochs,
            batch_size=args.batch_size,
            learning_rate=args.learning_rate
        )
        print(f"Dataset Comparison Results: {results}")
    elif args.command == "architecture-comparison":
        # Create architecture comparison experiment
        experiment = CrossArchitectureExperiment(
            experiment_manager=experiment_manager, 
            architectures=args.architectures, 
            dataset=args.dataset
        )
        
        # If output config is specified, save the configuration
        if hasattr(args, 'output_config') and args.output_config:
            # Create config
            config = ExperimentConfig(
                experiment_name=f"Architecture Comparison - {args.dataset}",
                dataset=args.dataset,
                model_architecture=args.architectures,
                epochs=args.epochs,
                batch_size=args.batch_size,
                learning_rate=args.learning_rate,
                cross_dataset_testing=False
            )
            
            # Save config file
            output_path = Path(args.output_config)
            if output_path.suffix.lower() in ['.yaml', '.yml']:
                config.save_yaml(output_path)
            else:
                config.save(output_path)
            print(f"Configuration saved to: {output_path}")
            return
        
        # Run experiment
        results = experiment.run(
            epochs=args.epochs,
            batch_size=args.batch_size,
            learning_rate=args.learning_rate
        )
        print(f"Architecture Comparison Results: {results}")
    elif args.command == "hyperparameter-optimization":
        # Create hyperparameter optimization experiment
        experiment = HyperparameterExperiment(
            experiment_manager=experiment_manager,
            model_architecture=args.architecture,
            dataset=args.dataset,
            n_trials=args.n_trials,
            timeout=args.timeout
        )
        
        # If output config is specified, save the configuration
        if hasattr(args, 'output_config') and args.output_config:
            # Create a base config for hyperparameter optimization
            config = ExperimentConfig(
                experiment_name=f"Hyperparameter Optimization - {args.architecture}",
                dataset=args.dataset,
                model_architecture=args.architecture,
                cross_dataset_testing=False
            )
            
            # Add hyperopt-specific settings
            hyperopt_config = {
                "experiment_type": "hyperparameter_optimization",
                "n_trials": args.n_trials,
                "timeout": args.timeout,
                "base_config": config.to_dict()
            }
            
            # Save config file
            output_path = Path(args.output_config)
            if output_path.suffix.lower() in ['.yaml', '.yml']:
                with open(output_path, 'w') as f:
                    yaml.dump(hyperopt_config, f, default_flow_style=False, sort_keys=False)
            else:
                with open(output_path, 'w') as f:
                    json.dump(hyperopt_config, f, indent=2)
            print(f"Configuration saved to: {output_path}")
            return
        
        # Run experiment
        results = experiment.run()
        print(f"Hyperparameter Optimization Results: {results}")
    elif args.command == "generate-report":
        results_compiler = ResultsCompiler()
        if args.format == "json":
            output_path = results_compiler.generate_comparative_report(args.experiment_ids, args.output)
            print(f"JSON report generated at: {output_path}")
        elif args.format == "yaml":
            # Generate report in YAML format
            report = results_compiler.compile_experiment_summary(args.experiment_ids[0]) if len(args.experiment_ids) == 1 else {
                "experiment_ids": args.experiment_ids,
                "summaries": [results_compiler.compile_experiment_summary(exp_id) for exp_id in args.experiment_ids]
            }
            
            output_path = args.output or Path(f"report_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.yaml")
            with open(output_path, 'w') as f:
                yaml.dump(report, f, default_flow_style=False, sort_keys=False)
            print(f"YAML report generated at: {output_path}")
        elif args.format == "pptx":
            output_path = results_compiler.generate_powerpoint_report(args.experiment_ids, args.output)
            print(f"PowerPoint report generated at: {output_path}")
        elif args.format == "excel":
            output_path = results_compiler.export_to_excel(args.experiment_ids, args.output)
            print(f"Excel report generated at: {output_path}")
    elif args.command == "resume":
        # Resume experiment using stored configuration
        results = experiment_manager.run_experiment(args.experiment_id)
        print(f"Resumed Experiment Results: {results}")
    else:
        parser.print_help()


if __name__ == "__main__":
    main() 